#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation for V2X DoS Vulnerabilities Research Report
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation (16:9 aspect ratio)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme (professional dark blue theme)
PRIMARY_COLOR = RGBColor(0x1a, 0x23, 0x7e)  # Dark blue
ACCENT_COLOR = RGBColor(0xe6, 0x8a, 0x00)   # Orange accent
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0x4a, 0x90, 0xd9)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Authors
    authors_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = authors_box.text_frame
    tf.text = "Siva Sai Akunuru | Sai Teja Duggelwar | Ravi Prakash Gudipudi | Sri Venkata Naga Sai Chinmai Malladi"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
    p.alignment = PP_ALIGN.CENTER
    
    # Course info
    course_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.5))
    tf = course_box.text_frame
    tf.text = "CYB402 – Vehicle to Everything (V2X) Cybersecurity"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_title):
    """Add section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = section_title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, subtitle=None):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle if provided
    y_pos = Inches(1.5)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.4))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = ACCENT_COLOR
        y_pos = Inches(1.9)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add two-column content slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    tf.text = left_title
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    tf.text = right_title
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    return slide

# ===== BUILD PRESENTATION =====

# Slide 1: Title
add_title_slide(prs, 
    "Automated Discovery of DoS Vulnerabilities in V2X Protocols",
    "A Critical Review of CVAnalyzer")

# Slide 2: Agenda
add_content_slide(prs, "Presentation Outline", [
    "Introduction to V2X Communication and Security Challenges",
    "Overview of the Assigned Research Paper (CVAnalyzer)",
    "System Model and Threat Assumptions",
    "Methodology: TLC Model Checking + PRISM Probabilistic Analysis",
    "Key Technical Contributions: State Space Reduction Techniques",
    "Evaluation Results: 19 Discovered DoS Vulnerabilities",
    "Critical Evaluation: Strengths and Limitations",
    "Real-World Implications for V2X Deployment",
    "Future Research Directions",
    "Conclusion and Q&A"
])

# Slide 3: Introduction - V2X
add_content_slide(prs, "Introduction: Vehicle-to-Everything (V2X) Communication", [
    "V2X enables real-time interaction between vehicles, infrastructure, pedestrians, and networks",
    "Built on DSRC and WAVE standards (IEEE 1609 family)",
    "Promises enhanced road safety, improved traffic efficiency, reduced environmental impact",
    "Uses PKI-based security with digital signatures for message integrity",
    "Critical safety systems rely on timely message delivery (<100ms latency per SAE J2945/1)"
])

# Slide 4: The Security Challenge
add_content_slide(prs, "The Security Challenge: DoS Threats in V2X", [
    "Wireless connectivity expands attack surface significantly",
    "Denial-of-Service (DoS) attacks pose severe risks to availability of safety-critical systems",
    "Potential consequences: blocked emergency notifications, disrupted traffic coordination, platoon failures",
    "Prior work relied on manual analysis and ad-hoc testing—insufficient for complex protocol stacks",
    "Need for systematic, automated vulnerability discovery before large-scale deployment"
])

# Slide 5: Paper Overview
add_content_slide(prs, "Paper Overview: CVAnalyzer", [
    "Paper: 'Automated Discovery of Denial-of-Service Vulnerabilities in Connected Vehicle Protocols'",
    "Authors: Hu et al. | Published: USENIX Security Symposium 2021",
    "CVAnalyzer combines TLC model checker (TLA+ suite) with PRISM probabilistic model checker",
    "Discovers design-level DoS vulnerabilities and quantifies their security/safety impact",
    "Target domains: IEEE 1609.2 P2PCD protocol and Platoon Management Protocols (VENTOS/PLEXE)"
])

# Slide 6: Key Contributions
add_content_slide(prs, "Key Contributions of the Research", [
    "Novel CVAnalyzer framework for systematic vulnerability discovery using soundness-focused verification",
    "Discovery of 19 DoS vulnerabilities: 4 in IEEE 1609.2 P2PCD (N1-N4) and 15 in platoon management (14 new)",
    "Quantitative assessment: N1 attack achieves 99% success rate; N2-N4 range from 41-87%",
    "All attacks violate SAE J2945/1 requirement of 100ms end-to-end latency",
    "Real-world testbed validation and responsible disclosure to IEEE 1609 Working Group"
])

# Slide 7: System Model
add_two_column_slide(prs, "System Model and Threat Assumptions",
    "System Scope",
    [
        "Complete CV communication stack modeled",
        "Focus on P2PCD protocol behavior",
        "Legitimate vehicle behavior and certificate patterns",
        "Protocol state machines encoded in TLA+",
        "Network: reliable FIFO channels assumed"
    ],
    "Threat Model",
    [
        "Attacker has legitimate credentials and certificates",
        "Can send cryptographically valid messages",
        "Goal: Cause DoS by exploiting design flaws",
        "Operates within protocol parameters (well-formed messages)",
        "Parameterized verification: n = 2, 3, 4 vehicles"
    ]
)

# Slide 8: The P2PCD Vulnerability
add_content_slide(prs, "The Critical Design Flaw: IEEE 1609.2 P2PCD", [
    "P2PCD (Peer-to-Peer Certificate Distribution) enables certificate sharing between vehicles",
    "Learning responses are sent as PDUs (Protocol Data Units), NOT SPDUs (Secured PDUs)",
    "PDUs lack digital signatures—cannot be cryptographically verified by recipient",
    "This unsigned nature enables attackers to inject forged learning responses",
    "Fundamental vulnerability in the protocol design affecting availability"
])

# Slide 9: Methodology Phase 1
add_content_slide(prs, "Methodology: Phase 1 – TLC Model Checking", [
    "TLC model checker from TLA+ suite exhaustively explores protocol state space",
    "Protocol behavior, attacker actions, and DoS properties encoded as TLA+ (~500 lines for P2PCD)",
    "TLC verifies whether DoS states are reachable under various attack scenarios",
    "Soundness prioritized over completeness—no false positives guaranteed",
    "Parameterized verification problem is undecidable; small vehicle numbers used (n=2,3,4)"
])

# Slide 10: State Space Reduction
add_two_column_slide(prs, "State Space Reduction Techniques",
    "Timer Index Reduction",
    [
        "Original hash space: M = 2^24 possible timer indices",
        "Reduced to range [0, X-1] where X = 5",
        "Dramatic reduction in state space explosion",
        "Maintains soundness—no false positives"
    ],
    "Event Equivalence Class Abstraction",
    [
        "Original event space: N = 2^256 possible events",
        "Reduced to [0, 2M-1] equivalence classes",
        "Maximum 9 classes when M=5",
        "Sound abstraction preserves attack validity"
    ]
)

# Slide 11: Methodology Phase 2
add_content_slide(prs, "Methodology: Phase 2 – PRISM Probabilistic Model Checking", [
    "PRISM quantifies likelihood and impact of discovered vulnerabilities",
    "Model checking alone cannot quantify probabilistic non-deterministic attacks (N1-N4)",
    "PMC assigns uniform probabilities to concurrent transitions",
    "Calculates attack success rates and expected latency impact",
    "Provides actionable metrics for prioritizing vulnerability remediation"
])

# Slide 12: Evaluation Results P2PCD
add_content_slide(prs, "Evaluation Results: IEEE 1609.2 P2PCD Protocol", [
    "4 DoS vulnerabilities discovered (N1-N4) exploiting unsigned learning responses",
    "N1-N2: Learning request flooding attacks",
    "N3-N4: Learning response manipulation attacks",
    "Attack success rates: N1 up to 99% under optimal timing; N2-N4 range 41-87%",
    "All four attacks violate SAE J2945/1 100ms latency requirement for safety-critical communications",
    "Non-deterministic attacks exploit protocol race conditions"
])

# Slide 13: Evaluation Results Platooning
add_content_slide(prs, "Evaluation Results: Platoon Management Protocols", [
    "Analysis of VENTOS and PLEXE platooning systems",
    "15 DoS vulnerabilities identified (14 are new discoveries)",
    "Includes deterministic attacks D1-D4 targeting coordination mechanisms",
    "Attacks force emergency braking, platoon dissolution, or hazardous spacing",
    "Vulnerabilities target safe following distances and synchronized speeds",
    "Diversity of vulnerabilities suggests systemic issues in platooning protocol design"
])

# Slide 14: Validation
add_content_slide(prs, "Real-World Validation and Impact", [
    "All 19 discovered attacks validated in real-world testbed with actual CV hardware/software",
    "Theoretical vulnerabilities confirmed as practical exploits",
    "Testbed measurements provided ground truth for probabilistic model predictions",
    "Resource consumption and latency impact measured empirically",
    "Findings reported to IEEE 1609 Working Group through responsible disclosure",
    "IEEE acknowledged vulnerabilities and plans to adopt proposed solutions"
])

# Slide 15: Strengths
add_content_slide(prs, "Critical Evaluation: Strengths", [
    "Systematic combination of TLC model checking + PRISM probabilistic analysis",
    "Soundness guarantee—no false positives through careful abstraction",
    "Sophisticated handling of state explosion via Timer Index Reduction and Event Equivalence Classes",
    "Real-world testbed validation strengthens practical relevance",
    "Quantified impact provides actionable guidance for risk prioritization",
    "Responsible disclosure and engagement with standardization bodies",
    "Methodology applicable to other resource-constrained, latency-sensitive systems"
])

# Slide 16: Limitations
add_content_slide(prs, "Critical Evaluation: Limitations", [
    "Parameterized verification limited to small vehicle numbers (n=2,3,4) due to state explosion",
    "Network abstraction assumes reliable FIFO channels—excludes packet loss and variable latency",
    "SCMS scope excludes certificate revocation from model",
    "99% success rate assumption may be difficult in high-relative-velocity scenarios (100+ km/h)",
    "State space reductions are sound but not complete—some vulnerabilities may be missed",
    "Scope limited to DoS only—excludes integrity, privacy, authentication bypass",
    "Focus on specific protocol versions—subsequent revisions may have addressed issues"
])

# Slide 17: Real-World Implications
add_content_slide(prs, "Real-World Implications for V2X Deployment", [
    "High success rates and severe impact suggest current protocols insufficient for safety-critical use",
    "SAE J2945/1 100ms latency violations compromise real-time safety guarantees",
    "Transportation authorities and manufacturers must consider risks in large-scale deployments",
    "Time lag between discovery and standard updates creates exposure window",
    "Defense-in-depth strategies needed: rate limiting, response validation, anomaly monitoring",
    "Highlights tension between performance optimization and security robustness"
])

# Slide 18: Future Research
add_content_slide(prs, "Future Research Directions", [
    "Extend CVAnalyzer to C-V2X (Cellular V2X) protocols based on LTE and 5G",
    "Develop techniques to scale parameterized verification to larger vehicle numbers",
    "Research compositional verification methods for realistic traffic densities",
    "Integrate machine learning-based anomaly detection with formal verification",
    "Automated patch generation for discovered vulnerabilities",
    "Expand analysis to privacy-preserving V2X technologies (pseudonym changes, anonymous credentials)",
    "Address intersection of security, privacy, and availability requirements"
])

# Slide 19: Conclusion
add_content_slide(prs, "Conclusion", [
    "CVAnalyzer demonstrates value of combining TLC model checking with PRISM probabilistic analysis",
    "19 DoS vulnerabilities discovered (4 P2PCD + 15 platooning, 14 new)",
    "Critical finding: P2PCD learning responses transmitted as unsigned PDUs",
    "Attack success rates up to 99% with latency violations of SAE J2945/1 requirements",
    "Sound abstraction techniques (Timer Index Reduction, Event Equivalence Classes) enable tractable verification",
    "Responsible disclosure practices positively influence industry standards",
    "Foundation for securing connected vehicle ecosystems as deployment scales"
])

# Slide 20: References/Q&A
add_content_slide(prs, "Key References", [
    "[1] Hu et al. 'Automated Discovery of DoS Vulnerabilities in Connected Vehicle Protocols.' USENIX Security 2021",
    "[2] IEEE Std 1609.2-2022: WAVE Security Services for Applications and Management Messages",
    "[3] SAE Std J2735: DSRC Message Set Dictionary",
    "[4] SAE Std J2945/1: On-Board System Requirements for V2V Safety Communications",
    "[5] PRISM 4.0: Verification of Probabilistic Real-Time Systems (CAV 2011)",
    "[6] TLA+ Home Page: https://lamport.azurewebsites.net/tla/tla.html",
    "[9] PLEXE: A Platooning Extension for Veins (IEEE VNC 2014)"
])

# Slide 21: Q&A
add_section_slide(prs, "Questions & Discussion")

# Save presentation
output_path = "/home/siva/.openclaw/workspace/V2X_DoS_Vulnerabilities_Presentation.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
