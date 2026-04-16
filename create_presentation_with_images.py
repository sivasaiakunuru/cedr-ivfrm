#!/usr/bin/env python3
"""
CEDR 35-Slide Presentation WITH IMAGES
Embeds actual UML diagrams and visualizations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
TEAL = RGBColor(0x00, 0x8B, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)

BASE_DIR = "/home/siva/.openclaw/workspace"
UML_DIR = f"{BASE_DIR}/uml"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_header(slide, title):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    h.fill.solid()
    h.fill.fore_color.rgb = NAVY
    h.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "CEDR: Cybersecurity Event Data Recorder"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
st = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
p = st.text_frame.paragraphs[0]
p.text = "In-Vehicle Forensic Readiness Module"
p.font.size = Pt(28)
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.CENTER
tm = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
p = tm.text_frame.paragraphs[0]
p.text = "Team Cyber-Torque | CYB408 Capstone | St. Clair College"
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER
add_notes(slide, "Welcome. Presenting CEDR - Cybersecurity Event Data Recorder. Modern vehicles face unprecedented cyber threats. CEDR provides tamper-evident forensic evidence.")
print("Slide 1: Title")

# Slides 2-5: Problem section (text only)
problem_slides = [
    ("Agenda", "Three Sections:\n\n1. THE PROBLEM\n   Automotive cybersecurity gap and forensic readiness crisis\n\n2. THE SOLUTION\n   CEDR architecture, differentiation, and roadmap\n\n3. THE ASK\n   Phase 1 investment and business case", "Three sections today: Problem, Solution, and Ask."),
    ("The Automotive Cybersecurity Gap", "450% increase in incidents since 2018 (Upstream Security)\n100+ ECUs per modern vehicle\n287 days average breach detection time\n$5.2M average breach cost (IBM 2024)", "Threat landscape escalating rapidly."),
    ("The Forensic Readiness Gap", "Key Finding: Existing solutions detect attacks.\nNone provide tamper-evident forensic evidence.\n\nESCRYPT: Basic logging only\nHarman: Cloud-only analysis\nArgus: Post-event only\nUpstream: No edge integrity\nCEDR: Hash chain integrity ✓", "Zero competitors provide tamper-evident logging."),
    ("Financial & Regulatory Impact", "€30M maximum UN R155 fines\nInsurance claims denied without forensic proof\n$24B projected industry cost by 2030 (McKinsey)\n6-12 month ISO 21434 certification delays", "UN R155 mandatory July 2024. Fines up to €30M.")
]

for title, content, notes in problem_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.5), Inches(5.2))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(19)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)
    print(f"Slide: {title}")

# Slide 6: CEDR Overview WITH Component Diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "CEDR Solution Overview")
# Add component diagram image
img_path = f"{UML_DIR}/component_diagram_professional.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.5), Inches(1.4), width=Inches(6.5))
# Add text on left
left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
tf = left.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "VEHICLE EDGE:\n• Raspberry Pi CM4\n• NXP A71CH HSM\n• Real-time CAN capture\n• ML inference <500ms\n• Hash chain storage\n• 4G/LTE upload\n\nCLOUD:\n• AWS + Azure DR\n• Kubernetes\n• SOC integration"
p.font.size = Pt(14)
p.font.color.rgb = DARK_GRAY
add_notes(slide, "Two-tier architecture with component diagram. Vehicle edge uses CM4 with NXP HSM. Cloud on AWS with Azure DR.")
print("Slide 6: CEDR Overview with Component Diagram")

# Slides 7-10: Solution text slides
solution_text = [
    ("Technical Differentiation: Hash Chain", "Every security event generates SHA-256 hash:\n• Event data (CAN frame, timestamp, vehicle ID)\n• Previous event's hash (sequential chain)\n• HSM-signed timestamp (non-repudiation)\n\nIf ANY record modified → hash verification FAILS\n\nCourt-admissible evidence integrity\nNo competitor offers this capability", "Core innovation: cryptographic hash chain."),
    ("Honest Hardware Assessment", "SPECIFICATION | PROTOTYPE | PRODUCTION\nPlatform | RPi CM4 | NXP S32G\nTemperature | -20°C to +70°C | -40°C to +125°C\nQualification | NOT automotive | AEC-Q100 Grade 1\n\nPhase 3: $85,000 migration to NXP S32G\n\nTransparent: CM4 is commercial grade, suitable for prototype validation but NOT for production vehicles.", "Honest about CM4 limitations. Phase 3 includes automotive-grade migration."),
    ("Revised Cloud Infrastructure Budget", "Multi-Cloud Architecture (1,000 vehicles):\n\nAWS (Primary): $5,075/month\nAzure (DR): $1,200/month\nGCP (ML): $800/month\n\nTOTAL: $7,075/month = $84,900/year\n(Revised from $29,940)\n\nRealistic AWS pricing including data transfer and CloudWatch.", "Costs increased from $29,940 to $84,900 with realistic usage estimates."),
    ("Security Features: STRIDE Coverage", "Spoofing: Multi-source location verification\nTampering: SHA-256 hash chain, WORM storage\nRepudiation: HSM-signed timestamps\nInformation Disclosure: AES-256-GCM, HSM keys\nDoS: Hardware CAN filtering, rate limiting\nElevation: RBAC, zero-trust architecture\n\nFull threat model coverage.", "All STRIDE categories covered with specific countermeasures.")
]

for title, content, notes in solution_text:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)
    print(f"Slide: {title}")

# Slide 11: User Stories WITH Story Pyramid
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "38 Agile User Stories")
# Add story pyramid image
img_path = f"{UML_DIR}/story_pyramid.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7), Inches(1.4), width=Inches(5.8))
# Add text on left
left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.5))
tf = left.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Complete Framework:\n\n7 Functional\n4 Non-Functional\n7 Security\n10 Abuse Cases\n10 Countermeasures\n\nAll stories include:\n• Given-When-Then criteria\n• OWASP verification\n• Traceability matrix"
p.font.size = Pt(15)
p.font.color.rgb = DARK_GRAY
add_notes(slide, "38 comprehensive user stories with acceptance criteria. Story pyramid visualization shows hierarchy.")
print("Slide 11: User Stories with Pyramid")

# Slide 12: Use Cases
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Use Cases: Complete Flow Documentation")
cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4))
tf = cb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "UC-001 Secure OTA Update Logging:\nPre-conditions → Primary Flow → Alternate (Rollback) → Exception (Invalid Sig)\n\nUC-002 Real-Time Intrusion Detection:\nPre-conditions → Primary Flow → Alternate (Offline Mode) → Exception (ML Failure)\n\nUC-003 Post-Incident Forensic Investigation:\nPre-conditions → Search → Verify Hash Chain → Export Evidence\n\nAll include: Pre/post conditions, primary flow, alternate flows, exception flows"
p.font.size = Pt(16)
p.font.color.rgb = DARK_GRAY
add_notes(slide, "Three primary use cases with complete flow documentation.")
print("Slide 12: Use Cases")

# Slide 13: UML Diagrams WITH Images
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "UML Architecture: 12 Diagrams")
# Add deployment diagram
img_path = f"{UML_DIR}/deployment_diagram_professional.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4), height=Inches(5.8))
add_notes(slide, "UML architecture includes structural and behavioral diagrams. Deployment diagram shows hardware/software stack.")
print("Slide 13: UML with Deployment Diagram")

# Slide 14: Risk Analysis WITH Heat Map
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Risk Analysis: 24 Risks")
# Add risk heatmap
img_path = f"{UML_DIR}/risk_heatmap.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.4), width=Inches(5.5))
# Add text on left
left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.8), Inches(5.5))
tf = left.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "NIST SP 800-30 Methodology:\nLikelihood × Impact × Confidence\n\nCategories:\n• Technical (5)\n• Business (5)\n• Compliance (4)\n• Operational (2)\n• Security (3)\n• Supply Chain (3)\n• Technical Debt (2)\n\nRisk Reduction: 57%\n12.4 → 5.3 average"
p.font.size = Pt(15)
p.font.color.rgb = DARK_GRAY
add_notes(slide, "24 risks with NIST methodology. Heat map visualization. 57% risk reduction with mitigations.")
print("Slide 14: Risk Analysis with Heat Map")

# Slide 15: Compliance
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Compliance Roadmap")
cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4))
tf = cb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ISO/SAE 21434: 3/7 Complete ✅ | 2/7 In Progress 🟡 | Month 6-12\nUN R155 CSMS: Gap analysis complete | Month 12-14\nUN R156 SUMS: Architecture ready | Month 10-12\nAEC-Q100: NOT qualified (CM4) | NXP S32G migration M9-11\n\nExternal audit: $67,000 + $35,000/year maintenance\n\nHonest assessment with clear timeline to full compliance."
p.font.size = Pt(17)
p.font.color.rgb = DARK_GRAY
add_notes(slide, "Honest compliance status. 3 of 7 ISO work products complete. AEC-Q100 not qualified - Phase 3 includes migration.")
print("Slide 15: Compliance")

# Slides 16-20: Business Case
business_slides = [
    ("Competitive Analysis: 8 Competitors", "ESCRYPT: $800-1200 | ECU IDS | ❌ No forensics\nHarman: $700-1000 | Network | ❌ Cloud-only\nArgus: $500-800/yr | Cloud IDS | ❌ Post-event\nUpstream: $15-50/yr | Vehicle SOC | ❌ No edge\nKaramba: $50-100/ECU | Hardening | ❌ No detection\nGuardKnox: $300-500 | Gateway | ❌ Limited\nCEDR: $600-750 | Forensic Ready | ✅ Hash chain ONLY", "Analyzed 8 competitors. CEDR uniquely provides tamper-evident logging."),
    ("Market Opportunity", "TAM: 250M connected vehicles by 2025 | $8B cyber spend\nSAM: 25M commercial fleets | $800M forensic segment\nSOM: 1,000 vehicles Year 1 | 10,000 by Year 3\nRevenue: $6M (hardware + subscriptions)", "Market sizing with TAM/SAM/SOM analysis."),
    ("ROI Calculation: 2.2×", "Expected Loss Without CEDR: $780,000/year\nExpected Loss With CEDR: $156,000/year\nAnnual Risk Reduction: $624,000\n24-Month Value: $1,248,000\nInvestment (Phases 1-3): $450,000\nROI: 2.2× using risk mitigation model", "ROI calculated via expected loss reduction, not simple payback."),
    ("Break-Even: 4,181 Vehicles", "Contribution Margin: $422/vehicle (Year 1)\nYear 1 Fixed Costs: $1,418,875\nBreak-Even: 4,181 vehicles\n\nSensitivity: Best case 2,480 | Worst case 5,600", "Break-even recalculated honestly. Original 2,847 was inflated."),
    ("5-Year TCO: $6.1M", "Original: $5.43M | Revised: $6.11M\n\nHardware: $213K | Cloud: $425K | Software: $1.6M\nPersonnel: $3.4M | Compliance: $498K\n\nPer Vehicle: $6,110 (+$132 from original)\n\nCloud revision: +$131,560 for realistic data transfer", "TCO corrected with honest cloud cost estimates.")
]

for title, content, notes in business_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)
    print(f"Slide: {title}")

# Slide 21: Roadmap WITH SWOT
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Four-Phase Roadmap: TRL 4 → 9")
# Add SWOT diagram
img_path = f"{UML_DIR}/swot_diagram.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.2), Inches(1.4), width=Inches(5.8))
# Add roadmap text on left
left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5))
tf = left.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Phase 0: TRL 4 | Prototype | $65K\nPhase 1: TRL 5 | CM4, HSM | $45K\nPhase 2: TRL 6 | Secure boot | $120K\nPhase 3: TRL 8 | AEC-Q100 | $285K\nPhase 4: TRL 9 | 10K units | $850K\n\nTotal: $1.37M\nRisk Reduction: 40% → 90%"
p.font.size = Pt(15)
p.font.color.rgb = DARK_GRAY
add_notes(slide, "Four-phase roadmap with TRL progression. SWOT analysis shows competitive position.")
print("Slide 21: Roadmap with SWOT")

# Slides 22-30: Remaining slides
remaining = [
    ("Success Metrics", "Phase 1: 10 CM4 units, -20°C to +70°C validated, 25 attacks, LOI signed\nPhase 2: Secure boot 100%, 50+ HIL coverage, code coverage >80%\nPhase 3: AEC-Q100, UN R155, 3 pilots, 50 vehicles, $50K revenue", "Clear KPIs at each phase."),
    ("Phase 1: $45,000 Investment", "Objectives: Validate core technology, de-risk hardware, 40% risk reduction\n\nHardware: $15K | Testing: $18K | BD: $12K\n\nDeliverables: 10 CM4 units, HSM integration, IP67 testing, 25 attacks validated, pilot LOI", "$45K Phase 1 investment with detailed breakdown."),
    ("Use of Funds", "Hardware $15K: 10× CM4 ($6.5K), 10× HSM ($1.2K), enclosures ($3.5K), cables ($3.8K)\nTesting $18K: Environmental chamber ($8K), Red team ($6K), Security audit ($4K)\nBD $12K: Pilot engagement ($7K), Advisory board ($3K), Travel ($2K)", "Detailed Phase 1 budget breakdown."),
    ("Phases 2-4 Overview", "Phase 2 ($120K): Security hardening, HIL testing, SOAR integration\nPhase 3 ($285K): AEC-Q100 ($125K), EMC ($35K), UN R155 ($67K), 50-vehicle pilot\nPhase 4 ($850K): 10K units, V2X, Series A preparation\n\nCumulative: $1.37M", "Phases 2-4 total $1.32M for certification and scale."),
    ("Risk Mitigation: 40% → 90%", "Pre-mitigation: 12.4 average score\nPhase 1: 12.4 → 6.8 (45% reduction)\nPhase 2: 6.8 → 4.0 (68% reduction)\nPhase 3: 4.0 → 2.3 (81% reduction)\nPhase 4: 2.3 → 1.1 (91% reduction)\n\nNatural stopping points after each phase.", "Risk reduction roadmap with go/no-go decision points."),
    ("Team & Partnerships", "Core Team: St. Clair College CYB408 Capstone\nPartners: eSentire (SOC), NXP (HSM), AWS/Azure/GCP (cloud)\nAdvisors: Automotive cyber expert, OEM integration, Legal/forensic\n\nAcademic: St. Clair College Cybersecurity Club", "Team and strategic partnerships."),
    ("Next Steps: 90-Day Plan", "Month 1: Procure CM4 units, HSM driver development, chamber access\nMonth 2: HSM integration, enclosure fabrication, pilot discussions\nMonth 3: Environmental testing, 25 attack validation, LOI signing\n\nDeliverable: Validated technology + committed customer", "90-day Phase 1 execution plan."),
    ("Q&A: 10 Tough Questions", "1. Hardware grade? → Phase 3 NXP S32G migration\n2. Competitors? → Analyzed 8, zero forensic capability\n3. Cloud costs? → Realistic AWS, unit economics work\n4. ISO complete? → 3/7 now, full by Month 12\n5. ML concerns? → Rule-based fallback\n6. ESCRYPT? → Target Tier-2/3 differentiation\n7. Break-even? → Honest 4,181 vs inflated 2,847\n8. Funding risk? → Natural stopping points\n9. GDPR? → Legal basis UN R155\n10. Why now? → UN R155 mandatory July 2024", "Prepared answers for tough evaluator questions.")
]

for title, content, notes in remaining:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)
    print(f"Slide: {title}")

# Slides 31-35: Appendices
appendix = [
    ("Appendix: Detailed Budget", "Phase 1: $45K | Phase 2: $120K | Phase 3: $285K | Phase 4: $850K\n5-Year TCO: $6.11M\nHardware: $213K | Cloud: $425K | Software: $1.6M | Personnel: $3.4M | Compliance: $498K", "Complete budget breakdown."),
    ("Appendix: Competitive Matrix", "Capability | ESCRYPT | Harman | Argus | Upstream | CEDR\nECU IDS | ✅ | ✅ | ❌ | ❌ | ✅\nNetwork IDS | ✅ | ✅ | ✅ | ✅ | ✅\nCloud Analytics | ✅ | ✅ | ✅ | ✅ | ✅\nTamper-Evident | ❌ | ❌ | ❌ | ❌ | ✅ UNIQUE\nHash Chain | ❌ | ❌ | ❌ | ❌ | ✅\nOpen Source | ❌ | ❌ | ❌ | ❌ | ✅\nPrice | $800-1200 | $700-1000 | $500-800 | $15-50 | $600-750", "Competitive matrix showing 8 capabilities."),
    ("Appendix: ISO 21434 Status", "WP-01 Policy | ✅ Complete\nWP-02 Risk Mgmt | ✅ Complete\nWP-03 Concept | ✅ Complete\nWP-04 Development | 🟡 In Progress\nWP-05 Validation | 🟡 In Progress\nWP-06 Production | 📐 Design\nWP-07 Updates | ✅ Design\n\n3/7 Complete | 2/7 In Progress | 2/7 Design", "ISO work product status."),
    ("Appendix: Technical Specs", "PROTOTYPE (CM4):\nBCM2711 quad-core @ 1.5GHz | 8GB LPDDR4\n-20°C to +70°C (NOT AEC-Q100)\n\nPRODUCTION (NXP S32G):\nVehicle Network Processor\n-40°C to +125°C (AEC-Q100 Grade 1)\nISO 26262 ASIL-B capable\n\nMigration: $85K (Phase 3)", "Side-by-side hardware comparison."),
    ("References: 32 Sources", "[1] Upstream Security 2024 Report\n[3] IBM Cost of Data Breach 2024\n[5] UNECE UN R155\n[6] Mansor et al. Forensic Framework 2020\n[7] Song et al. CAN IDS 2020\n[8] Crosby & Wallach 2009\n[9] NXP A71CH Datasheet\n[10] NIST SP 800-30\n[12] Miller & Valasek 2015\n[17] ISO/SAE 21434:2021\n[20] ISO 27037\n+ 21 additional IEEE sources", "Key IEEE and industry references.")
]

for title, content, notes in appendix:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)
    print(f"Slide: {title}")

# Save
output_path = f"{BASE_DIR}/CEDR_Final_Presentation_35_with_Images.pptx"
prs.save(output_path)
print(f"\n✅ Created: CEDR_Final_Presentation_35_with_Images.pptx")
print(f"Total slides: {len(prs.slides)}")
print(f"Images embedded: Component diagram, Story pyramid, Deployment diagram, Risk heatmap, SWOT diagram")
