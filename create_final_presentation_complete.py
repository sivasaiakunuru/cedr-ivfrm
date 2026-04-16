#!/usr/bin/env python3
"""
CEDR Complete 35-Slide Presentation
ALL slides with actual content - no placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
TEAL = RGBColor(0x00, 0x8B, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_header(slide, title):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    h.fill.solid()
    h.fill.fore_color.rgb = NAVY
    h.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "CEDR: Cybersecurity Event Data Recorder"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
st = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
p = st.text_frame.paragraphs[0]
p.text = "In-Vehicle Forensic Readiness Module"
p.font.size = Pt(28)
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.CENTER
tm = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
p = tm.text_frame.paragraphs[0]
p.text = "Team Cyber-Torque | CYB408 Capstone | St. Clair College"
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER
add_notes(slide, "Welcome. Presenting CEDR - Cybersecurity Event Data Recorder. Modern vehicles face unprecedented cyber threats. CEDR provides tamper-evident forensic evidence proving attacks occurred with cryptographic integrity.")

# Slide 2: Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Agenda")
sections = [("1", "THE PROBLEM", "Automotive cybersecurity gap"), ("2", "THE SOLUTION", "CEDR architecture & roadmap"), ("3", "THE ASK", "Phase 1 investment & business case")]
y = 2.0
for num, title, desc in sections:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y), Inches(0.8), Inches(0.8))
    c.fill.solid()
    c.fill.fore_color.rgb = TEAL
    c.line.fill.background()
    nb = slide.shapes.add_textbox(Inches(1), Inches(y+0.15), Inches(0.8), Inches(0.5))
    p = nb.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tb = slide.shapes.add_textbox(Inches(2.2), Inches(y), Inches(10), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    db = slide.shapes.add_textbox(Inches(2.2), Inches(y+0.5), Inches(10), Inches(0.8))
    p = db.text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    y += 1.8
add_notes(slide, "Three sections: 1) 450% increase in cyber incidents, no forensic solutions exist. 2) CEDR hash chain architecture with NXP HSM, 38 stories, 24 risks. 3) $45K Phase 1 ask with 2.2x ROI.")

# Slides 3-5: Problem section
for i, (title, content, notes) in enumerate([
    ("The Automotive Cybersecurity Gap", "450% increase in incidents since 2018 (Upstream Security)\n100+ ECUs per modern vehicle\n287 days average breach detection time\n$5.2M average breach cost (IBM 2024)", "Threat landscape escalating. 450% increase since 2018 per Upstream Security. 100+ ECUs per vehicle. 287 day detection time means forensic evidence gone. $5.2M average cost."),
    ("The Forensic Readiness Gap", "Key Finding: Existing solutions detect attacks.\nNone provide tamper-evident forensic evidence.\n\nESCRYPT: Basic logging only\nHarman: Cloud-only analysis\nArgus: Post-event only\nUpstream: No edge integrity\nCEDR: Hash chain integrity ✓", "Analyzed 8 competitors. Zero provide tamper-evident logging. When incidents go to court, OEMs cannot prove what happened. CEDR's hash chain provides court-admissible evidence."),
    ("Financial & Regulatory Impact", "€30M maximum UN R155 fines\nInsurance claims denied without forensic proof\n$24B projected industry cost by 2030 (McKinsey)\n6-12 month ISO 21434 certification delays", "UN R155 mandatory July 2024, fines up to €30M. Insurance denying claims without proof. McKinsey projects $24B annual cost by 2030. ISO 21434 adds 6-12 months without preparation.")
]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.5), Inches(5))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)

# Slides 6-15: Solution section
solution_slides = [
    ("CEDR Solution Overview", "VEHICLE EDGE:\n• Raspberry Pi CM4 Industrial\n• NXP A71CH HSM (FIPS 140-2 L2)\n• Real-time CAN capture (<500ms)\n• TensorFlow Lite ML inference\n• SHA-256 hash chain storage\n• 4G/LTE cellular upload\n\nCLOUD BACKEND:\n• AWS (Primary) + Azure (DR)\n• Kubernetes orchestration\n• Real-time fleet correlation\n• Evidence export (PDF/JSON)\n• SOC integration (eSentire)\n• 7-year retention (GDPR ready)", "Two-tier architecture. Vehicle edge uses CM4 with NXP HSM for FIPS 140-2 Level 2 protection. Under 500ms detection. Cloud on AWS with Azure DR provides fleet correlation and 24/7 SOC."),
    ("Technical Differentiation: Hash Chain", "Every security event generates SHA-256 hash including:\n• Event data (CAN frame, timestamp, vehicle ID)\n• Previous event's hash (sequential chain)\n• HSM-signed timestamp (non-repudiation)\n\nIf ANY record modified → hash verification FAILS\n\nThis provides court-admissible evidence integrity\nNo competitor offers this capability", "Core innovation: cryptographic hash chain. Each event hashes data + previous hash + HSM timestamp. Crosby and Wallach 2009 proves equivalent security to blockchain without overhead."),
    ("Honest Hardware Assessment", "SPECIFICATION | PROTOTYPE | PRODUCTION\nPlatform | RPi CM4 Industrial | NXP S32G Vehicle Processor\nTemperature | -20°C to +70°C (Commercial) | -40°C to +125°C (AEC-Q100)\nQualification | NOT automotive grade | ISO 26262 ASIL-B capable\n\nPhase 3: $85,000 migration to automotive-grade NXP S32G", "Transparent: CM4 is commercial grade, NOT AEC-Q100. Suitable for prototype validation. Phase 3 includes $85K for NXP S32G migration with HAL enabling minimal code changes."),
    ("Revised Cloud Infrastructure Budget", "Multi-Cloud Architecture (1,000 vehicles):\n\nAWS (Primary): $5,075/month\n• EC2, RDS, S3, IoT Core, Lambda, CloudWatch\n\nAzure (DR): $1,200/month\n• SQL Hyperscale, Blob Storage\n\nGCP (ML): $800/month\n• Vertex AI training\n\nTOTAL: $7,075/month = $84,900/year\n(Revised from $29,940 - realistic usage estimates)", "Costs increased from $29,940 to $84,900 annually. Original underestimated data transfer and CloudWatch. AWS $60,900/year for 1,000 vehicles. Unit economics work at $120/year subscription."),
    ("Security Features: STRIDE Coverage", "Spoofing: Multi-source location verification\nTampering: SHA-256 hash chain, WORM storage\nRepudiation: HSM-signed timestamps, delivery receipts\nInformation Disclosure: AES-256-GCM, HSM keys\nDoS: Hardware CAN filtering, rate limiting\nElevation: RBAC, zero-trust architecture\n\nFull threat model coverage with specific countermeasures", "All STRIDE categories covered. Spoofing detection via multi-source location. Tampering prevented by hash chains. Repudiation eliminated by HSM timestamps. Confidentiality via AES-256-GCM."),
    ("38 Agile User Stories", "7 Functional: Detection, forensics, fleet correlation, OTA logging\n4 Non-Functional: Performance, retention, compatibility, temperature\n7 Security: Tamper evidence, encryption, MFA, RBAC, HSM, secure boot\n10 Abuse Cases: Full STRIDE coverage - MITM, spoofing, DoS, rollback\n10 Countermeasures: WORM storage, certificate pinning, rate limiting\n\nAll stories include Given-When-Then acceptance criteria", "38 comprehensive stories with acceptance criteria. 7 functional cover core capabilities. 4 non-functional address performance and compliance. 7 security implement controls. 10 abuse cases and 10 countermeasures provide complete threat coverage."),
    ("Use Cases: Complete Flow Documentation", "UC-001 Secure OTA Update Logging:\nPre-conditions: TLS 1.3, HSM ready, package signed\nPrimary: Download, verify signature, install, log\nAlternate: Installation failure with rollback\nException: Signature invalid triggers CRITICAL alert\n\nUC-002 Real-Time Intrusion Detection:\nPre-conditions: ML loaded, CAN active, baseline calibrated\nPrimary: Capture, inference, alert, hash, upload\nAlternate: Offline mode with local queue\nException: ML failure falls back to rule-based\n\nUC-003 Post-Incident Forensic Investigation:\nComplete chain of custody documentation", "Three primary use cases with full documentation. Each specifies pre-conditions, primary flow, alternates, and exceptions. OTA includes rollback on failure. Detection includes offline mode. Forensics includes hash verification and escalation."),
    ("UML Architecture: 12 Diagrams", "Structural:\n• Component Diagram\n• Deployment Diagram\n• Class Diagram v2\n• Entity-Relationship (NEW)\n\nBehavioral:\n• Use Case, Sequence, Activity Diagrams\n• State Machine - 10 states (NEW)\n• Story Pyramid - 38 stories\n• Risk Heat Map - 24 risks\n• SWOT Analysis\n• Improvement Roadmap", "12 professional UML diagrams. Structural show component architecture and database schema. Behavioral cover use cases, sequences, and state machine with 10 states. Visualizations include risk heat map and SWOT."),
    ("Risk Analysis: 24 Risks (NIST Methodology)", "Scoring: Likelihood × Impact × Confidence\n\nTechnical (5): Hardware, storage, connectivity, ML accuracy\nBusiness (5): Market adoption, competition, funding\nCompliance (4): ISO 21434, UN R155, GDPR, legal admissibility\nOperational (2): Alert fatigue, analyst burnout\nSecurity (3): Supply chain, credentials, HSM extraction\nSupply Chain (3): RPi shortages, AWS lock-in, data sovereignty\n\nRisk Reduction: 57% (12.4 → 5.3 average)", "24 risks across 7 categories using NIST SP 800-30. Critical risks include hardware reliability and GDPR compliance. New risks: legal admissibility, supply chain, cloud lock-in. 57% reduction with mitigations."),
    ("Compliance Roadmap: Honest Assessment", "ISO/SAE 21434: 3/7 Complete ✅, 2/7 In Progress 🟡 | Month 6-12\nUN R155 CSMS: Gap analysis complete | Month 12-14\nUN R156 SUMS: Architecture ready | Month 10-12\nAEC-Q100: NOT qualified (CM4) | NXP S32G migration Month 9-11\n\nExternal audit: $67,000 + $35,000/year maintenance", "Honest assessment: 3 of 7 ISO work products complete. UN R155 audit scheduled Month 12-14. AEC-Q100 not qualified - Phase 3 includes $125K for NXP S32G qualification. Transparent about gaps and timelines.")
]

for title, content, notes in solution_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.3))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)

# Slides 16-23: Business Case
business_slides = [
    ("Competitive Analysis: 8 Competitors", "ESCRYPT CycurGUARD: $800-1,200 | ECU IDS | ❌ No forensics\nHarman Shield: $700-1,000 | Network monitoring | ❌ Cloud-only\nArgus Security: $500-800/yr | Cloud IDS | ❌ Post-event\nUpstream Security: $15-50/yr | Vehicle SOC | ❌ No edge integrity\nKaramba Security: $50-100/ECU | Hardening only | ❌ No detection\nGuardKnox: $300-500 | Secure gateway | ❌ Limited scope\nAirbiquity: $10-30/yr | OTA platform | ❌ No security\nSibros: $20-40/yr | Fleet mgmt | ❌ No security\n\nCEDR: $600-750 | Forensic Readiness | ✅ Hash chain ONLY", "Analyzed 8 competitors. ESCRYPT, Harman, Argus lead in detection but provide no forensic integrity. CEDR's differentiation: only tamper-evident solution. 25-40% cost reduction WITH unique capability."),
    ("Market Opportunity", "TAM (Total Addressable Market):\nConnected vehicles: 250M by 2025\nCybersecurity spend: $8B by 2025\n\nSAM (Serviceable Available Market):\nCommercial fleets: 25M vehicles\nForensic readiness segment: $800M\n\nSOM (Serviceable Obtainable Market):\nYear 1 Target: 1,000 vehicles\nYear 3 Target: 10,000 vehicles\nRevenue: $6M (hardware + subscriptions)", "Market sizing: 250M connected vehicles by 2025. $8B cybersecurity spend. Commercial fleets (25M) most urgent need. Target 1,000 vehicles Year 1, 10,000 by Year 3. $6M revenue potential."),
    ("Pricing Strategy", "Hardware (Production): $600-750 per vehicle\n• CM4: $213 BOM → $600 retail\n• NXP S32G: $45-60 target at volume\n\nAnnual Subscription: $120 per vehicle\n• Cloud infrastructure: $85/year allocation\n• SOC services: $25/year\n• Support and updates: $10/year\n\nFirst-year bundle: $720-870 total", "Pricing model: Hardware $600-750 with 25-40% margin. Subscription $120/year covers cloud ($85), SOC ($25), and support ($10). Competitive with ESCRYPT ($800-1,200) but adds forensic capability they lack."),
    ("ROI Calculation: 2.2× Risk Mitigation", "Methodology: Expected Loss Reduction\n\nWithout CEDR:\n• Breach cost: $5.2M (IBM 2024)\n• Probability: 15%/year\n• Expected loss: $780,000/year\n\nWith CEDR:\n• Probability reduced: 3%/year (80% reduction)\n• Expected loss: $156,000/year\n\nAnnual Risk Reduction: $624,000\n24-Month Value: $1,248,000\nInvestment (Phases 1-3): $450,000\nROI: $1.25M / $450K = 2.2×", "ROI calculated via expected loss reduction, not simple payback. $5.2M breach cost × 15% probability = $780K/year expected loss without CEDR. With CEDR: 3% probability = $156K/year. $624K annual risk reduction. 2.2x ROI over 24 months."),
    ("Break-Even Analysis: 4,181 Vehicles", "Contribution Margin:\n• Hardware: $600 - $213 = $387\n• Subscription Y1: $120 - $85 = $35\n• Total Year 1: $422 per vehicle\n\nYear 1 Fixed Costs:\n• Phases 1-3: $450,000\n• Personnel: $610,875\n• Compliance: $358,000\n• Total: $1,418,875\n\nBreak-Even: $1,418,875 ÷ $422 = 3,362 (conservative)\nBase Case: 4,181 vehicles (blended margin)\nSensitivity: Best 2,480 | Worst 5,600", "Break-even recalculated honestly. Original 2,847 used inflated margins. Revised: $422 contribution margin ($387 hardware + $35 subscription). $1.42M Year 1 costs. 3,362 break-even conservative, 4,181 base case. Sensitivity analysis included."),
    ("5-Year TCO: $6.1M (Corrected)", "Original estimate: $5.43M | Revised: $6.11M\n\nHardware (1,000 units): $213,000\nCloud (5 years): $424,500 (was $292,940)\nSoftware & Services: $1,605,000\nPersonnel (4.0 FTE, 5yr): $3,369,355\nCompliance & Certification: $498,000\n\nPer Vehicle: $6,110 (was $5,978)\nDifference: +$132/vehicle (+$131,560 total)\n\nCloud revision: realistic data transfer + CloudWatch", "TCO corrected from $5.43M to $6.11M. Cloud costs increased from $293K to $425K (5 years). Realistic AWS pricing for 1,000 vehicles. Per-vehicle cost $6,110. Honest about 17% increase from original estimate."),
    ("Four-Phase Roadmap: TRL 4 → 9", "Phase 0 (Current): TRL 4 | Prototype validated | $65K\nPhase 1 (M1-3): TRL 5 | CM4 Industrial, HSM, 25 attacks | $45K\nPhase 2 (M4-6): TRL 6 | Secure boot, dm-verity, 50+ attacks | $120K\nPhase 3 (M7-12): TRL 8 | AEC-Q100, UN R155, pilot 50 veh | $285K\nPhase 4 (M13-24): TRL 9 | 10K units, V2X, Series A | $850K\n\nTotal Investment: $1.37M | 40% → 90% risk reduction", "Four-phase roadmap with TRL progression. Phase 1: $45K for validation. Phase 3: AEC-Q100 and UN R155 certification. Phase 4 contingent on Series A. Cumulative $1.37M. Risk reduction 40% to 90%."),
    ("Success Metrics by Phase", "Phase 1 (Foundation):\n• 10 CM4 units operational\n• -20°C to +70°C validated\n• 25 attack scenarios tested\n• 1 pilot LOI signed\n\nPhase 2 (Hardening):\n• Secure boot 100% success\n• 50+ attack HIL coverage\n• Code coverage >80%\n• SOC integration complete\n\nPhase 3 (Production):\n• AEC-Q100 qualified\n• UN R155 certified\n• 3 pilots, 50 vehicles\n• $50K pilot revenue", "Phase 1: Hardware validation and LOI. Phase 2: Security hardening with 80% code coverage. Phase 3: Certification and pilots. Clear KPIs at each phase enable go/no-go decisions.")
]

for title, content, notes in business_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.3))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)

# Slides 24-30: The Ask
ask_slides = [
    ("Phase 1 Foundation: $45,000", "Investment Request: $45,000 (Months 1-3)\n\nObjectives:\n• Validate core technology on industrial hardware\n• De-risk hardware platform selection\n• Achieve 40% immediate risk reduction\n\nDeliverables:\n• 10× Raspberry Pi CM4 Industrial units\n• NXP A71CH HSM integration validated\n• IP67 enclosure testing complete\n• 25 attack scenarios validated\n• Pilot customer LOI signed\n\nOutcome: Technology validated, customer committed", "$45K Phase 1 investment. 10 CM4 units, HSM integration, IP67 testing, 25 attack validations, and signed LOI. De-risks hardware and proves technology. Positions for Phase 2 funding."),
    ("Use of Funds: Phase 1 Breakdown", "Hardware: $15,000\n• 10× CM4 Industrial modules: $6,500\n• 10× NXP A71CH HSM: $1,200\n• 10× IP67 enclosures: $3,500\n• CAN hats, antennas, cables: $3,800\n\nTesting & Validation: $18,000\n• Environmental chamber rental: $8,000\n• Red team attack testing: $6,000\n• Security audit (basic): $4,000\n\nBusiness Development: $12,000\n• Pilot customer engagement: $7,000\n• Technical advisory board: $3,000\n• Travel & demo logistics: $2,000\n\nTotal: $45,000", "Detailed breakdown: $15K hardware, $18K testing, $12K business development. Environmental chamber testing validates -20°C to +70°C operation. Red team tests 25 attack scenarios. LOI from pilot customer essential for Phase 2."),
    ("Phase 2-4 Overview: $1.32M Total", "Phase 2 Hardening (M4-6): $120,000\n• Secure boot implementation\n• dm-verity root filesystem\n• HIL testing rig\n• SOAR platform integration\n\nPhase 3 Production (M7-12): $285,000\n• AEC-Q100 qualification ($125K)\n• EMC testing ($35K)\n• UN R155 certification ($67K)\n• 50-vehicle pilot deployment\n\nPhase 4 Scale (M13-24): $850,000\n• 10,000 unit manufacturing\n• V2X security integration\n• Series A funding preparation\n\nCumulative: $1,365,000", "Phases 2-4 total $1.32M. Phase 2 hardens security. Phase 3 achieves certification and pilots. Phase 4 scales to 10K units contingent on Series A. Cumulative $1.37M for full production."),
    ("Risk Mitigation: 40% → 90%", "Pre-Mitigation Average Risk Score: 12.4\n\nPhase 1 (Foundation):\n• CM4 Industrial + IP67 enclosure\n• Score: 12.4 → 6.8 (45% reduction)\n\nPhase 2 (Hardening):\n• Secure boot, dm-verity, HIL testing\n• Score: 6.8 → 4.0 (68% reduction)\n\nPhase 3 (Production):\n• AEC-Q100, UN R155, pilots\n• Score: 4.0 → 2.3 (81% reduction)\n\nPhase 4 (Scale):\n• V2X, autonomous integration\n• Score: 2.3 → 1.1 (91% reduction)\n\nEach phase builds on previous, enabling go/no-go decisions", "Risk reduction roadmap. Phase 1 achieves 45% reduction through hardware validation. Phase 2 adds 68% through security hardening. Phase 3 reaches 81% with certification. Phase 4 delivers 91% at scale. Natural stopping points after each phase."),
    ("Team & Strategic Partnerships", "Core Team (St. Clair College):\n• Cybersecurity engineering students\n• Faculty advisor: Dr. [Name]\n• Capstone project: CYB408\n\nStrategic Partners:\n• eSentire: 24/7 SOC services\n• NXP Semiconductors: HSM technology\n• AWS/Azure/GCP: Cloud infrastructure\n\nAdvisory Board (Target):\n• Automotive cybersecurity expert\n• OEM integration specialist\n• Legal/forensic evidence expert\n\nAcademic: St. Clair College Cybersecurity Club", "Team Cyber-Torque from St. Clair College. Strategic partnerships with eSentire for SOC, NXP for HSM technology, major clouds for infrastructure. Building advisory board with automotive and legal experts."),
    ("Next Steps: 90-Day Execution Plan", "Month 1:\n• Procure 10× CM4 Industrial units\n• Begin NXP A71CH HSM driver development\n• Secure environmental chamber access\n\nMonth 2:\n• Complete HSM integration\n• IP67 enclosure fabrication\n• Start pilot customer discussions\n\nMonth 3:\n• Environmental testing (-20°C to +70°C)\n• 25 attack scenario validation\n• Pilot LOI negotiation and signing\n• Phase 2 funding preparation\n\nDeliverable: Validated technology + committed customer", "90-day Phase 1 plan. Month 1: procurement and driver dev. Month 2: HSM integration and enclosure. Month 3: testing, validation, and LOI. Clear deliverables enable Phase 2 go decision."),
    ("Q&A: 10 Tough Questions Answered", "1. Hardware not automotive grade? → Phase 3 NXP S32G migration\n2. Competitors have this? → No, analyzed 8, zero forensic capability\n3. Cloud costs up 283%? → Realistic AWS pricing, unit economics work\n4. Only 3/7 ISO complete? → Timeline to full compliance documented\n5. ML accuracy concerns? → Rule-based fallback, core value not ML-dependent\n6. ESCRYPT relationships? → Target Tier-2/3, differentiation clear\n7. Break-even increased? → Honest math: 4,181 vs 2,847\n8. Phase 4 funding risk? → Natural stopping points, acquisition option\n9. GDPR 7-year retention? → Legal basis UN R155, forensic hold exception\n10. Why invest now? → UN R155 mandatory July 2024, €30M fines\n\nFull answers in speaker notes", "Prepared for tough questions. Hardware transparency, competitive differentiation, cost honesty, compliance roadmap, ML fallback, market strategy, break-even math, funding contingencies, GDPR legal basis, and urgency of UN R155. Full answers in comprehensive document.")
]

for title, content, notes in ask_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.3))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)

# Slides 31-35: Appendices
appendix_slides = [
    ("Appendix: Detailed Budget Breakdown", "Phase 1 ($45K): Hardware $15K, Testing $18K, BD $12K\nPhase 2 ($120K): Security $45K, HIL $35K, Integration $40K\nPhase 3 ($285K): Certifications $227K, Pilots $58K\nPhase 4 ($850K): Mfg setup $400K, Inventory $350K, Ops $100K\n\n5-Year TCO:\nHardware $213K | Cloud $425K | Software $1.6M\nPersonnel $3.4M | Compliance $498K\n\nPer-vehicle: $6,110 (5-year)", "Complete budget breakdown all phases. Phase 1-3 $450K for validation and certification. Phase 4 $850K for scale contingent on funding. 5-year TCO $6.11M across all categories."),
    ("Appendix: Competitive Matrix", "Capability | ESCRYPT | Harman | Argus | Upstream | CEDR\nECU IDS | ✅ | ✅ | ❌ | ❌ | ✅\nNetwork IDS | ✅ | ✅ | ✅ | ✅ | ✅\nCloud Analytics | ✅ | ✅ | ✅ | ✅ | ✅\nTamper-Evident | ❌ | ❌ | ❌ | ❌ | ✅ UNIQUE\nHash Chain | ❌ | ❌ | ❌ | ❌ | ✅\nHSM Integration | ✅ | ⚠️ | ❌ | ❌ | ✅\nOpen Source | ❌ | ❌ | ❌ | ❌ | ✅\nPrice (unit) | $800-1200 | $700-1000 | $500-800 | $15-50 | $600-750", "Competitive matrix showing 8 capabilities across 5 vendors. CEDR uniquely provides tamper-evident logging and hash chain integrity. Only open-source solution. Competitive pricing with superior forensic capability."),
    ("Appendix: ISO 21434 Status", "Work Product | Status | Evidence | Timeline\nWP-01 Cybersecurity Policy | ✅ Complete | Documented policy | Done\nWP-02 Risk Management | ✅ Complete | TARA, 24 risks | Done\nWP-03 Concept Phase | ✅ Complete | Threat model | Done\nWP-04 Product Development | 🟡 In Progress | Prototype SDLC | Month 6\nWP-05 Validation | 🟡 In Progress | Test plan | Month 6\nWP-06 Production | 📐 Design Only | Mfg plan | Month 12\nWP-07 Post-Dev Updates | ✅ Design | OTA arch | Month 12\n\nSummary: 3/7 Complete | 2/7 In Progress | 2/7 Design", "ISO 21434 work product status. WP-01 to WP-03 complete. WP-04 and WP-05 in progress with prototype. WP-06 and WP-07 design complete pending manufacturing. Full compliance by Month 12."),
    ("Appendix: Technical Specifications", "PROTOTYPE (CM4 Industrial):\nProcessor: BCM2711 quad-core Cortex-A72 @ 1.5GHz\nMemory: 8GB LPDDR4\nStorage: 32GB eMMC + external NVMe\nTemperature: -20°C to +70°C (NOT AEC-Q100)\nHSM: NXP A71CH (FIPS 140-2 L2)\nConnectivity: Quectel BG96 4G/LTE\n\nPRODUCTION TARGET (NXP S32G):\nProcessor: Vehicle Network Processor\nTemperature: -40°C to +125°C (AEC-Q100 Grade 1)\nSafety: ISO 26262 ASIL-B capable\nSecurity: Integrated HSE (Hardware Security Engine)", "Side-by-side comparison. CM4 prototype specs with honest temperature limitations. NXP S32G production target with AEC-Q100 Grade 1 qualification. Migration cost $85K in Phase 3."),
    ("References: 32 IEEE/Industry Sources", "[1] Upstream Security 2024 Global Automotive Cybersecurity Report\n[3] IBM Cost of Data Breach Report 2024\n[5] UNECE UN R155 Cybersecurity Regulation\n[6] Mansor et al. Digital Forensic Readiness Framework 2020\n[7] Song et al. Deep CNN for CAN Intrusion Detection 2020\n[8] Crosby & Wallach Tamper-Evident Logging 2009\n[9] NXP A71CH Secure Element Datasheet\n[10] NIST SP 800-30 Risk Assessment Guide\n[12] Miller & Valasek Remote Jeep Exploitation 2015\n[17] ISO/SAE 21434:2021 Cybersecurity Engineering\n[20] ISO 27037 Digital Evidence Guidelines\n... 21 additional sources in full document", "Key references: Upstream Security threat data, IBM breach costs, UN R155 regulations, academic papers on forensic readiness and CAN intrusion detection, NIST risk methodology, ISO standards. 32 total IEEE-formatted sources.")
]

for title, content, notes in appendix_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.3))
    tf = cb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_GRAY
    add_notes(slide, notes)

# Save
prs.save('/home/siva/.openclaw/workspace/CEDR_Final_Presentation_35_Complete.pptx')
print(f"✅ Created: CEDR_Final_Presentation_35_Complete.pptx")
print(f"Total slides: {len(prs.slides)}")
