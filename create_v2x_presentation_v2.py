#!/usr/bin/env python3
"""
Professional V2X DoS Vulnerabilities Presentation - Concise & Visual
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Professional color palette
DARK_BLUE = RGBColor(0x1a, 0x23, 0x7e)
ACCENT_ORANGE = RGBColor(0xe6, 0x8a, 0x00)
LIGHT_BLUE = RGBColor(0x2e, 0x5c, 0x9e)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xc4, 0x1e, 0x3a)
GREEN = RGBColor(0x28, 0x8e, 0x3c)

def add_title_slide(prs):
    """Title slide with gradient-like effect"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Decorative accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Automated Discovery of DoS Vulnerabilities in V2X Protocols"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    tf.text = "A Critical Review of CVAnalyzer"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Authors
    authors_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.333), Inches(1))
    tf = authors_box.text_frame
    authors = "Siva Sai Akunuru • Sai Teja Duggelwar • Ravi Prakash Gudipudi • Sri Venkata Naga Sai Chinmai Malladi"
    tf.text = authors
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xbb, 0xbb, 0xbb)
    p.alignment = PP_ALIGN.CENTER
    
    # Course
    course_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4))
    tf = course_box.text_frame
    tf.text = "CYB402 – V2X Cybersecurity"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_header(prs, title):
    """Section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BLUE
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, highlight_key_points=True):
    """Content slide with professional layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header background
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12.133), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Check if this is a key stat point
        if highlight_key_points and any(x in bullet for x in ['19', '99%', '100ms', '14 new']):
            p.text = f"▸ {bullet}"
            p.font.color.rgb = RED
            p.font.bold = True
        else:
            p.text = f"▸ {bullet}"
            p.font.color.rgb = DARK_GRAY
        
        p.font.size = Pt(20)
        p.space_before = Pt(14)
        p.space_after = Pt(6)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    tf.text = left_title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    tf.text = right_title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.8), Inches(4.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)
    
    return slide

def add_visual_slide(prs, title, main_points, visual_type="architecture"):
    """Slide with visual diagram representation using shapes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if visual_type == "cvanalyzer":
        # Draw CVAnalyzer architecture diagram
        # Input box
        inp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2), Inches(2.5), Inches(0.8))
        inp.fill.solid()
        inp.fill.fore_color.rgb = LIGHT_GRAY
        inp.line.color.rgb = LIGHT_BLUE
        inp.line.width = Pt(2)
        tf = inp.text_frame
        tf.text = "Protocol\nSpecifications"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # TLC box
        tlc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(2), Inches(2.5), Inches(0.8))
        tlc.fill.solid()
        tlc.fill.fore_color.rgb = RGBColor(0xe3, 0xf2, 0xfd)
        tlc.line.color.rgb = LIGHT_BLUE
        tlc.line.width = Pt(2)
        tf = tlc.text_frame
        tf.text = "TLC Model\nChecker"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
        
        # PRISM box
        prism = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2), Inches(2.5), Inches(0.8))
        prism.fill.solid()
        prism.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)
        prism.line.color.rgb = GREEN
        prism.line.width = Pt(2)
        tf = prism.text_frame
        tf.text = "PRISM\nProbabilistic MC"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
        
        # Output box
        out = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(2), Inches(2.5), Inches(0.8))
        out.fill.solid()
        out.fill.fore_color.rgb = RGBColor(0xff, 0xeb, 0xee)
        out.line.color.rgb = RED
        out.line.width = Pt(2)
        tf = out.text_frame
        tf.text = "Vulnerabilities\n+ Quantified Risk"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RED
        
        # Arrows (lines)
        for x in [3.3, 6.3, 9.3]:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.35), Inches(0.5), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_ORANGE
            line.line.fill.background()
        
        # Key points below
        y_start = Inches(3.3)
    else:
        y_start = Inches(1.7)
    
    # Content points
    content_box = slide.shapes.add_textbox(Inches(0.6), y_start, Inches(12.133), Inches(3.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(main_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if any(x in point for x in ['99%', '19', '100ms', '14 new']):
            p.text = f"▸ {point}"
            p.font.color.rgb = RED
            p.font.bold = True
        else:
            p.text = f"▸ {point}"
            p.font.color.rgb = DARK_GRAY
        
        p.font.size = Pt(18)
        p.space_before = Pt(12)
    
    return slide

def add_results_slide(prs):
    """Special slide for results with visual stats"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Key Results: 19 DoS Vulnerabilities Discovered"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Stat boxes
    # Box 1: Total vulnerabilities
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(2.8), Inches(1.6))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RED
    box1.line.fill.background()
    tf = box1.text_frame
    tf.text = "19\nTotal\nVulnerabilities"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = WHITE
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    tf.paragraphs[2].font.size = Pt(14)
    tf.paragraphs[2].font.color.rgb = WHITE
    tf.paragraphs[2].alignment = PP_ALIGN.CENTER
    
    # Box 2: P2PCD
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(1.7), Inches(2.8), Inches(1.6))
    box2.fill.solid()
    box2.fill.fore_color.rgb = ACCENT_ORANGE
    box2.line.fill.background()
    tf = box2.text_frame
    tf.text = "4\nP2PCD\nAttacks (N1-N4)"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = WHITE
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    tf.paragraphs[2].font.size = Pt(14)
    tf.paragraphs[2].font.color.rgb = WHITE
    tf.paragraphs[2].alignment = PP_ALIGN.CENTER
    
    # Box 3: Platooning
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.7), Inches(2.8), Inches(1.6))
    box3.fill.solid()
    box3.fill.fore_color.rgb = GREEN
    box3.line.fill.background()
    tf = box3.text_frame
    tf.text = "15\nPlatooning\n(14 new)"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = WHITE
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    tf.paragraphs[2].font.size = Pt(14)
    tf.paragraphs[2].font.color.rgb = WHITE
    tf.paragraphs[2].alignment = PP_ALIGN.CENTER
    
    # Box 4: Success rate
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(1.7), Inches(2.5), Inches(1.6))
    box4.fill.solid()
    box4.fill.fore_color.rgb = LIGHT_BLUE
    box4.line.fill.background()
    tf = box4.text_frame
    tf.text = "99%\nN1 Attack\nSuccess Rate"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(12)
    tf.paragraphs[1].font.color.rgb = WHITE
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    tf.paragraphs[2].font.size = Pt(12)
    tf.paragraphs[2].font.color.rgb = WHITE
    tf.paragraphs[2].alignment = PP_ALIGN.CENTER
    
    # Details below
    details_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12.133), Inches(3.5))
    tf = details_box.text_frame
    tf.word_wrap = True
    
    points = [
        "All attacks violate SAE J2945/1 requirement of 100ms end-to-end latency",
        "N1-N4 exploit unsigned P2PCD learning responses (PDUs vs SPDUs)",
        "Platooning attacks (D1-D4) cause emergency braking and platoon dissolution",
        "All vulnerabilities validated in real-world testbed with actual CV hardware"
    ]
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(14)
    
    return slide

def add_qa_slide(prs):
    """Final Q&A slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Large Q&A text
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(2))
    tf = qa_box.text_frame
    tf.text = "Questions & Discussion"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1.5))
    tf = contact_box.text_frame
    tf.text = "Thank You"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ===== BUILD CONCISE PRESENTATION (12-15 slides) =====

# Slide 1: Title
add_title_slide(prs)

# Slide 2: The Problem
add_content_slide(prs, "The V2X Security Challenge", [
    "V2X enables real-time vehicle-to-vehicle and vehicle-to-infrastructure communication",
    "Built on DSRC/WAVE standards (IEEE 1609) with PKI-based security",
    "DoS attacks pose severe risks—compromising availability of safety-critical systems",
    "Prior manual analysis insufficient for complex protocol interactions",
    "Need: Systematic, automated vulnerability discovery before large-scale deployment"
])

# Slide 3: Paper Overview
add_visual_slide(prs, "CVAnalyzer: System Overview", [
    "Paper: Hu et al., USENIX Security 2021",
    "Combines TLC model checker (TLA+) with PRISM probabilistic model checking",
    "Discovers design-level DoS vulnerabilities + quantifies security/safety impact",
    "Targets: IEEE 1609.2 P2PCD protocol and VENTOS/PLEXE platooning"
], "cvanalyzer")

# Slide 4: The Core Vulnerability
add_content_slide(prs, "The Critical Design Flaw: P2PCD Protocol", [
    "Peer-to-Peer Certificate Distribution (P2PCD) enables certificate sharing between vehicles",
    "CRITICAL: Learning responses sent as PDUs (unsigned) instead of SPDUs (signed)",
    "PDUs lack digital signatures—cannot be cryptographically verified",
    "Attackers can inject forged learning responses without detection",
    "Fundamental protocol design flaw enabling non-deterministic attacks (N1-N4)"
])

# Slide 5: Methodology
add_two_column_slide(prs, "Methodology: Two-Phase Analysis",
    "Phase 1: TLC Model Checking",
    [
        "Exhaustive state space exploration",
        "Protocol encoded in TLA+ (~500 lines)",
        "Timer Index Reduction: 2²⁴ → 5 states",
        "Event Equivalence: 2²⁵⁶ → 9 classes",
        "Soundness guarantee—no false positives"
    ],
    "Phase 2: PRISM Analysis",
    [
        "Probabilistic model checking",
        "Quantifies non-deterministic attacks",
        "Calculates success probabilities",
        "Measures latency impact",
        "Validates real-world exploitability"
    ]
)

# Slide 6: Key Results (Visual stats)
add_results_slide(prs)

# Slide 7: Attack Details
add_content_slide(prs, "Attack Breakdown: P2PCD Vulnerabilities", [
    "N1-N2: Learning request flooding attacks—overwhelm certificate distribution",
    "N3-N4: Learning response manipulation—inject forged certificates",
    "Attack success rates: N1 achieves 99% under optimal timing conditions",
    "N2-N4 success rates: 41-87% depending on network conditions",
    "ALL attacks violate SAE J2945/1 100ms latency requirement"
])

# Slide 8: Platooning Attacks
add_content_slide(prs, "Platoon Management Protocol Vulnerabilities", [
    "15 vulnerabilities in VENTOS and PLEXE platooning systems (14 new discoveries)",
    "Deterministic attacks (D1-D4) target coordination mechanisms",
    "Attacks force: emergency braking events, platoon dissolution, hazardous spacing",
    "Target safe following distances and synchronized speeds at high velocity",
    "Diversity suggests systemic protocol design issues"
])

# Slide 9: Strengths
add_content_slide(prs, "Critical Evaluation: Strengths", [
    "Novel combination of TLC + PRISM for systematic vulnerability discovery",
    "Soundness guarantee—no false positives through careful abstraction",
    "Effective state explosion handling via Timer Index and Event Equivalence reductions",
    "Real-world testbed validation confirms practical exploitability",
    "Quantified impact enables risk prioritization",
    "Responsible disclosure to IEEE 1609 Working Group"
])

# Slide 10: Limitations
add_content_slide(prs, "Critical Evaluation: Limitations", [
    "Parameterized verification limited to n=2,3,4 vehicles (state explosion)",
    "Network model assumes reliable FIFO—excludes packet loss, variable latency",
    "SCMS scope excludes certificate revocation",
    "99% success rate may be difficult at high relative velocities (100+ km/h)",
    "Sound but not complete—some vulnerabilities may be missed",
    "Limited to DoS; excludes integrity, privacy, authentication bypass"
])

# Slide 11: Implications
add_content_slide(prs, "Real-World Implications", [
    "High success rates (up to 99%) and severe impact suggest protocols insufficient for safety-critical use",
    "SAE J2945/1 latency violations compromise real-time safety guarantees",
    "Time lag between discovery and standard updates creates exposure window",
    "Defense-in-depth needed: rate limiting, response validation, anomaly monitoring",
    "Highlights tension between performance optimization and security robustness"
])

# Slide 12: Future Work
add_content_slide(prs, "Future Research Directions", [
    "Extend to C-V2X (Cellular V2X) protocols using LTE/5G",
    "Scale parameterized verification to larger vehicle numbers",
    "Integrate ML-based anomaly detection with formal verification",
    "Automated patch generation for discovered vulnerabilities",
    "Analyze privacy-preserving V2X (pseudonym changes, anonymous credentials)"
])

# Slide 13: Conclusion
add_content_slide(prs, "Conclusion", [
    "CVAnalyzer successfully combines TLC model checking + PRISM probabilistic analysis",
    "19 DoS vulnerabilities discovered: 4 P2PCD + 15 platooning (14 new)",
    "Critical design flaw: P2PCD learning responses transmitted as unsigned PDUs",
    "Attack success up to 99% with violations of 100ms latency requirements",
    "Sound abstraction techniques enable tractable verification",
    "Responsible disclosure influences industry standards positively"
])

# Slide 14: References
add_content_slide(prs, "Key References", [
    "[1] Hu et al. 'Automated Discovery of DoS Vulnerabilities in Connected Vehicle Protocols.' USENIX Security 2021",
    "[2] IEEE Std 1609.2-2022: WAVE Security Services",
    "[4] SAE Std J2945/1: V2V Safety Communications Requirements",
    "[5] PRISM: Probabilistic Model Checker (CAV 2011)",
    "[6] TLA+ and TLC Model Checker (Lamport)"
], highlight_key_points=False)

# Slide 15: Q&A
add_qa_slide(prs)

# Save
output_path = "/home/siva/.openclaw/workspace/V2X_DoS_Vulnerabilities_Presentation.pptx"
prs.save(output_path)
print(f"✓ Professional presentation created: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
print(f"✓ Format: 16:9 widescreen, dark blue theme with visual elements")
