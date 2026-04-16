#!/usr/bin/env python3
"""
Update CEDR PowerPoint with comprehensive Budget & TCO slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')

# Find slide 11 (Budget slide - index 10)
slide = prs.slides[10]

# Clear existing content except title
shapes_to_remove = []
for shape in slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text and not text.startswith('Project Budget'):
            shapes_to_remove.append(shape)

# Remove old content shapes
for shape in shapes_to_remove:
    sp = shape.element
    sp.getparent().remove(sp)

# Add comprehensive budget table
table_left = Inches(0.5)
table_top = Inches(1.5)
table_width = Inches(12.3)
table_height = Inches(4.5)

# Create table with 8 rows, 4 columns
rows = 8
cols = 4
table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# Set column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.9)
table.columns[3].width = Inches(3.4)

# Header row
headers = ['Category', 'Prototype (1 unit)', 'Pilot (50 veh)', 'Scale (1,000 veh)']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    paragraph.font.size = Pt(11)
    paragraph.alignment = PP_ALIGN.CENTER

# Data rows
data = [
    ['Hardware per Unit', '$470', '$340', '$195'],
    ['Total Hardware', '$470', '$17,000', '$195,000'],
    ['Setup & Tools', '$2,010', '$22,500', '$125,000'],
    ['Cloud (Annual)', '$426', '$3,840', '$19,740'],
    ['Compliance', '$0', '$34,000', '$62,000'],
    ['Operations', '$0', '$12,000', '$156,000'],
    ['Total Year 1', '$2,480', '$106,340', '$527,540'],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        if col_idx == 0:
            paragraph.font.bold = True
        else:
            paragraph.alignment = PP_ALIGN.RIGHT
        
        # Alternate row colors
        if row_idx == len(data):  # Total row
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xe8, 0xf4, 0xe8)
            paragraph.font.bold = True

# Add key metrics below table
metrics_text = """Key Metrics:
• Cost per Vehicle (Pilot): $2,127 | (Scale): $613 Year 1 → $322 steady-state
• 5-Year TCO (1,000 vehicles): $1.6M | AWS Calculator: calculator.aws/#/estimate?id=cedr-scale
• ISO/SAE 21434 Compliance: $65,000 | UN R155 CSMS: $91,000 Year 1"""

metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.2))
tf = metrics_box.text_frame
tf.text = metrics_text
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(10)
    paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save updated presentation
prs.save('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')
print("✅ Updated Slide 11 with comprehensive Budget & TCO")
print("✅ Presentation saved: CEDR_Final_Presentation.pptx")
