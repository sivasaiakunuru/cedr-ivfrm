import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

DARK_NAVY = RGBColor(0x0A, 0x16, 0x28)
ELECTRIC_BLUE = RGBColor(0x00, 0xA8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Read the existing presentation
prs = Presentation('cdk_global_case_study.pptx')

# We want to replace the text boxes that have "[Image Placeholder]" with actual pictures
# Slide mapping to image files
slide_images = {
    0: "images/slide1.jpg",
    1: "images/slide2.jpg",
    2: "images/slide3.png", # Pie chart
    3: "images/slide4.jpg",
    4: "images/slide5.jpg",
    5: "images/slide6.jpg",
    6: "images/slide7.png", # Bar chart
    7: "images/slide8.jpg",
    8: "images/slide9.jpg",
    9: "images/slide10.jpg",
    10: "images/slide11.jpg",
    11: "images/slide12.jpg",
    12: "images/slide13.jpg",
    13: "images/slide14.jpg",
    14: "images/slide15.jpg",
    15: "images/slide16.jpg"
}

for idx, slide in enumerate(prs.slides):
    if idx in slide_images and os.path.exists(slide_images[idx]):
        img_path = slide_images[idx]
        
        # Determine placement
        if idx == 0:
            # Title slide background or center
            left = Inches(8.0)
            top = Inches(2.0)
            width = Inches(4.5)
        else:
            # Delete the image placeholder text box if found
            shapes_to_delete = []
            for shape in slide.shapes:
                if shape.has_text_frame and "[Image Placeholder]" in shape.text:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    shapes_to_delete.append(shape)
            
            # Note: deleting shapes in python-pptx can be tricky, we'll try to remove its element
            for shape in shapes_to_delete:
                sp = shape._element
                sp.getparent().remove(sp)

            # Fallback placement if no placeholder was found
            if not shapes_to_delete:
                left = Inches(8.0)
                top = Inches(1.5)
                width = Inches(4.8)

        # Insert picture
        try:
            slide.shapes.add_picture(img_path, left, top, width=width)
        except Exception as e:
            print(f"Failed to add image to slide {idx+1}: {e}")

prs.save("cdk_global_case_study_with_images.pptx")
print("Presentation updated with images.")
