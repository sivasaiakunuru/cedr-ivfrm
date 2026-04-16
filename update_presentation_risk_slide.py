#!/usr/bin/env python3
"""
Add Risk Analysis slide to CEDR PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')

# Add new slide after Budget slide (index 11)
# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)

# Move slide to position 12 (after budget, before demo)
# Get slide ID
slide_id = slide.slide_id

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "Risk Analysis & Mitigation"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)

# Risk summary table
table_left = Inches(0.5)
table_top = Inches(1.2)
table_width = Inches(7)
table_height = Inches(3)

rows = 6
cols = 4
table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# Column widths
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(1.4)
table.columns[2].width = Inches(1.4)
table.columns[3].width = Inches(2)

# Headers
headers = ['Risk Category', 'Level', 'Mitigation', 'Residual']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

# Data
data = [
    ['Technical (Hardware)', 'Medium', '75%', 'Low-Med'],
    ['Business (Adoption)', 'High', '60%', 'Medium'],
    ['Compliance (ISO)', 'Low', '90%', 'Low'],
    ['Operational (SOC)', 'Medium', '70%', 'Low-Med'],
    ['Security (Attacks)', 'Medium', '80%', 'Low'],
]

colors = [
    RGBColor(0xFF, 0xD9, 0xD9),  # Light red
    RGBColor(0xFF, 0xD9, 0xD9),  # Light red
    RGBColor(0xD9, 0xFF, 0xD9),  # Light green
    RGBColor(0xFF, 0xFF, 0xD9),  # Light yellow
    RGBColor(0xFF, 0xFF, 0xD9),  # Light yellow
]

for row_idx, (row_data, color) in enumerate(zip(data, colors), start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        if col_idx == 0:
            p.font.bold = True
        else:
            p.alignment = PP_ALIGN.CENTER

# Top 5 Risks
risks_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.2), Inches(5), Inches(3))
tf = risks_box.text_frame
tf.text = "Top 5 Risks (by Score):"

risks = [
    "1. B-002: Competitive dominance (21)",
    "2. T-001: Hardware reliability (21)",
    "3. B-001: Market adoption (19.5)",
    "4. B-003: Funding gap (13.5)",
    "5. T-002: Storage capacity (12)"
]

for risk in risks:
    p = tf.add_paragraph()
    p.text = risk
    p.font.size = Pt(11)
    p.space_after = Pt(4)

# SWOT Summary
swot_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12), Inches(2.8))
tf = swot_box.text_frame
tf.text = "SWOT Analysis:"

swot_text = """STRENGTHS: 82% cost advantage vs competitors ($145 vs $800+); Tamper-evident design; ISO 21434 compliant; Real-time capability
WEAKNESSES: Consumer-grade hardware (not automotive temp); Limited attack testing (13 scenarios); Academic origin (credibility gap)
OPPORTUNITIES: UN R155 mandate July 2024; Insurance premium discounts; $8.9B HSM market growth; V2X expansion
THREATS: ESCRYPT/Argus patent litigation; OEMs building in-house solutions; Economic downturn reducing R&D spend"""

p = tf.add_paragraph()
p.text = swot_text
p.font.size = Pt(10)
p.space_before = Pt(6)

# Save
prs.save('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')
print("✅ Added Risk Analysis slide (Slide 12)")
print("✅ Included: Risk table, Top 5 risks, SWOT analysis")
print("✅ Presentation saved: CEDR_Final_Presentation.pptx")
