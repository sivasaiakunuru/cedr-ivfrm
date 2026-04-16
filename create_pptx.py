import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_NAVY = RGBColor(0x0A, 0x16, 0x28)
    ELECTRIC_BLUE = RGBColor(0x00, 0xA8, 0xE8)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    slides_data = [
        {
            "title": "CDK Global Cyberattack (June 2024)",
            "subtitle": "Business Continuity Failure Across the North American Auto Dealership Ecosystem\nCYB 406 | St. Clair College | [Team Names]",
            "content": "A comprehensive analysis of the June 2024 cyberattack.",
            "notes": "Welcome everyone. Today we are presenting Case Study 2: The CDK Global Cyberattack of June 2024. This presentation addresses critical failures in Business Continuity and Disaster Recovery within the automobility ecosystem. We'll explore how a single point of failure at a major software provider crippled over 15,000 dealerships across North America, resulting in an estimated $1 billion in direct financial losses. This case directly connects to our CYB 406 course learning outcomes, specifically how we develop, analyze, and critique DR and BC plans in an automotive context. Let's start by looking at our roadmap for today's discussion. [Transition: Let's review our agenda.]",
            "image_ref": "Resources #1, #5 (CDK & St. Clair Logos) | #2, #4 (Background Concept)"
        },
        {
            "title": "Agenda / Presentation Roadmap",
            "content": "1. Introduction to CDK Global\n2. Incident Timeline\n3. Threat Actor Profile (BlackSuit)\n4. BIA: Operational & Financial Impact\n5. BC/DR Evaluation (RTO/RPO/MAD)\n6. Root Cause Analysis\n7. NIST CSF 2.0 Mapping\n8. Legal, Regulatory & Ethical Issues\n9. Recommendations & Case Study Comparison\n10. Key Takeaways & Q&A",
            "notes": "Here is our roadmap for the next 15-20 minutes. We will begin with an introduction to CDK Global to understand their scale and why they are a systemic risk. We'll trace the incident timeline, paying special attention to the critical failure point: the premature restoration attempt. Then, we will profile the BlackSuit ransomware gang, map the operational and financial impact using Business Impact Analysis frameworks, and critically evaluate the expected versus actual Recovery Time Objectives (RTO). We will also map deficiencies against NIST CSF 2.0, review legal/ethical issues, compare this to our previous C-Commerce case study, and close with actionable recommendations. [Transition: First, let's understand who CDK Global is.]",
            "image_ref": "Resources #6, #7, #8 (Roadmap Infographic & Icons)"
        },
        {
            "title": "Introduction to CDK Global",
            "content": "• Major SaaS provider for auto dealerships\n• Serves ~15,000 dealerships in North America\n• Core product: Dealer Management System (DMS)\n• Controls ~40% market share (DMS oligopoly)\n• Critical hub for sales, financing, and service operations\n• Ecosystem interdependency creates high systemic risk",
            "notes": "To understand the impact of the cyberattack, we must understand CDK Global's footprint. CDK is a leading SaaS provider whose Dealer Management System (DMS) operates as the central nervous system for approximately 15,000 auto dealerships. Together with Reynolds & Reynolds, they control roughly 70% of the market. The DMS integrates everything: inventory, payroll, financing, parts, and service. It is a critical hub. When a dealership loses its DMS, it reverts to the 1980s—using pen, paper, and calculators. This high level of integration means that CDK is a single point of failure for the entire automobility retail ecosystem. [Transition: Let's examine how the attack unfolded.]",
            "image_ref": "Resources #9, #10, #11, #12, #13 (Ecosystem Diagram, Market Share Chart)"
        },
        {
            "title": "The Attack: Incident Timeline (June - July 2024)",
            "content": "• June 18: Initial breach detected, systems proactively shut down\n• June 19: Premature restoration attempt\n• June 19 (Late): Second attack occurs during restoration\n• June 20-30: Complete blackout, ransom negotiation phase\n• Early July: Phased restoration begins\n• Late July: Most core DMS services finally restored",
            "notes": "The timeline of this incident reveals a textbook disaster recovery failure. On June 18, CDK detected a breach and shut down its systems to contain it. Dealerships were told systems would be back up shortly. On June 19, CDK attempted to restore services prematurely. This critical error exposed their recovery environment, leading to a second, devastating cyberattack by the BlackSuit ransomware gang later that evening. The systems were shut down again, initiating a blackout that lasted weeks. It wasn't until late July that most dealerships were fully back online. This highlights the danger of rushing the 'Recover' phase without proper containment and validation. [Transition: Who was behind this attack?]",
            "image_ref": "Resources #14-19 (Timeline Infographic & News Screenshots)"
        },
        {
            "title": "Threat Actor Profile: BlackSuit Ransomware",
            "content": "• Successor to Conti and Royal ransomware gangs\n• Employs double-extortion model (encrypt + steal)\n• Advanced TTPs mapped to MITRE ATT&CK\n• High initial access sophistication (phishing, compromised creds)\n• Targeted critical backup and recovery infrastructure\n• Demanded an estimated $25M+ ransom payout",
            "notes": "The threat actor responsible is BlackSuit, a sophisticated ransomware group believed to be the successor to the notorious Conti and Royal gangs. BlackSuit utilizes a double-extortion model: they don't just encrypt the data; they exfiltrate it and threaten to leak sensitive dealership and consumer financial records. Mapping their TTPs using the MITRE ATT&CK framework reveals highly sophisticated initial access and defense evasion techniques. In the CDK attack, BlackSuit explicitly targeted backup systems, eliminating CDK's ability to easily roll back to a clean state. Reports indicate a ransom of tens of millions was paid to facilitate decryption. [Transition: Let's look at the operational impact on dealerships.]",
            "image_ref": "Resources #20-25 (ATT&CK Map, Kill Chain, Ransomware Tree)"
        },
        {
            "title": "BIA: Operational Impact",
            "content": "• Sales: Handwritten contracts, unable to process financing\n• Service: Lost repair histories, manual parts ordering\n• Accounting: Payroll delays, cash flow interruption\n• Duration: ~14 to 30+ days of degraded operations\n• Forced transition to manual BC workarounds\n• Severe loss of operational visibility and customer trust",
            "notes": "A Business Impact Analysis (BIA) highlights massive operational disruption. For over two weeks, dealerships had to operate manually. In the Sales department, staff resorted to handwriting contracts and couldn't process loan approvals efficiently, killing deals. The Service department, a major revenue driver, couldn't access vehicle repair histories or electronically order parts. Technicians literally stood idle. Accounting departments struggled with payroll and floorplan financing. The extended duration—far beyond any reasonable Maximum Acceptable Outage (MAO)—pushed dealer manual workarounds to their breaking point, demonstrating that modern dealerships lack adequate offline contingency plans. [Transition: What was the financial cost of this operational halt?]",
            "image_ref": "Resources #26-30 (BIA Framework, Impact Heat Map, Handwritten forms)"
        },
        {
            "title": "BIA: Financial Impact",
            "content": "• Direct Dealer Losses: Estimated at $1B+ (Anderson Economic Group)\n• CDK Global Costs: Ransom payment ($25M estimated), remediation\n• Benchmark: Avg. ransomware breach cost is $5.13M (IBM 2024)\n• Secondary Impacts: Lost vehicle sales, delayed service revenue\n• Long-term Costs: Class-action lawsuits, regulatory fines\n• Stock/Market Impact: Immediate plunge in automotive retail stocks",
            "notes": "The financial impact was staggering. According to the Anderson Economic Group, direct losses to the dealership network exceeded $1 billion in just the first two weeks, primarily from lost vehicle sales and halted service bays. This far exceeds the IBM 2024 average ransomware breach cost of $5.13 million. For CDK Global, costs include the reported $25M ransom payment, emergency IT remediation, and upcoming legal liabilities. The incident also caused a temporary dip in the stock prices of major publicly traded dealership groups like AutoNation and Sonic Automotive. This clearly illustrates how cyber risk translates directly into systemic financial risk. [Transition: Let's evaluate this against standard recovery metrics.]",
            "image_ref": "Resources #31-35 (IBM Data Breach Chart, Financial Impact Graph)"
        },
        {
            "title": "BC/DR Evaluation: RTO, RPO & MAD",
            "content": "• RTO (Recovery Time Objective): Expected <24h vs. Actual 14-30 days\n• RPO (Recovery Point Objective): Data loss mitigation mostly successful\n• MAD (Maximum Acceptable Downtime): Exceeded (industry MAD ~3 days)\n• Critical Failure: Recovery environment compromise\n• Gap Analysis: Massive deviation from expected BCP capabilities\n• Result: Total failure of the 'Availability' security triad principle",
            "notes": "When we evaluate CDK's disaster recovery execution, we see a massive failure in meeting objectives. A typical Recovery Time Objective (RTO) for a tier-1 critical service like a DMS should be less than 24 hours. The actual recovery took 14 to 30 days. The Maximum Acceptable Downtime (MAD) for a dealership before severe financial distress occurs is generally around 3 to 5 days; this was vastly exceeded. While the Recovery Point Objective (RPO) seems to have been met—meaning massive data loss didn't permanently corrupt historical records—the failure to meet RTO was due to the compromise of the recovery environment during the premature restoration attempt. [Transition: Let's identify the root causes of these failures.]",
            "image_ref": "Resources #36-40 (RTO/RPO Diagram, ISO 22301 Lifecycle, Gap Analysis)"
        },
        {
            "title": "Root Cause Analysis",
            "content": "• Lack of adequate network segmentation\n• Single Point of Failure (SPOF) in cloud architecture\n• Inadequate isolated backups (Failure of 3-2-1 rule)\n• Premature system restoration without proper containment\n• Over-reliance on a single vendor by dealerships\n• Insufficient incident response testing and validation",
            "notes": "Using a root cause analysis, we identify six primary failures. First, inadequate network segmentation allowed the ransomware to spread from initial access to core infrastructure. Second, CDK acted as a massive Single Point of Failure. Third, the fact that BlackSuit delayed recovery suggests a failure of the 3-2-1 backup rule—CDK likely lacked immutable, air-gapped backups. Fourth, as noted, the premature restoration allowed the second attack. Fifth, the dealerships themselves failed by lacking redundant systems or robust offline procedures. Finally, this points to a systemic lack of rigorous, full-scale incident response and disaster recovery testing. [Transition: Let's map these deficiencies to the NIST framework.]",
            "image_ref": "Resources #41-45 (Fishbone Diagram, Network Segmentation, 3-2-1 Backup)"
        },
        {
            "title": "NIST CSF 2.0 Mapping",
            "content": "• GOVERN: Supply chain risk management failure\n• IDENTIFY: Underestimated systemic risk of monolithic architecture\n• PROTECT: Insufficient network segmentation and access controls\n• DETECT: Delayed anomaly detection in core systems\n• RESPOND: Flawed containment strategy (premature restoration)\n• RECOVER: Unacceptable RTO, compromised recovery environment",
            "notes": "If we map this incident to the newly updated NIST Cybersecurity Framework 2.0, CDK showed deficiencies across all six core functions. Under GOVERN, there was a failure in supply chain risk management. Under IDENTIFY, the risk of a monolithic architecture was underestimated. Under PROTECT, access controls and network segmentation failed. Under DETECT, the initial breach wasn't caught before lateral movement occurred. Under RESPOND, the containment strategy was fatally flawed, leading to the second attack. Finally, under RECOVER, the compromised backup environment led to an unacceptable RTO. This radar chart visualizes the massive gap between expected framework adherence and actual performance. [Transition: What are the legal implications?]",
            "image_ref": "Resources #46-49 (NIST CSF 2.0 Wheel, Radar Chart)"
        },
        {
            "title": "Legal & Regulatory Analysis",
            "content": "• FTC Safeguards Rule: Stricter reporting and security requirements\n• State Data Breach Notification Laws: 50-state compliance matrix\n• PIPEDA (Canada): Cross-border privacy and notification impacts\n• Class-Action Lawsuits: Multiple suits filed alleging negligence\n• SEC Reporting Requirements: Material cybersecurity incident disclosures\n• UNECE WP.29 / ISO 21434: Broader automotive supply chain context",
            "notes": "The legal fallout is highly complex. Dealerships and CDK are subject to the revised FTC Safeguards Rule, which mandates strict security programs and reporting for financial institutions—which dealerships are classified as. In Canada, PIPEDA governs the exposure of consumer data. Furthermore, CDK faces a myriad of class-action lawsuits alleging gross negligence in protecting sensitive consumer personally identifiable information (PII). Because CDK is deeply embedded in the automotive supply chain, this also touches upon broader regulatory themes seen in UNECE WP.29 and ISO/SAE 21434 regarding automotive cybersecurity lifecycle management. [Transition: Beyond the law, we must consider the ethics.]",
            "image_ref": "Resources #50-56 (FTC Rule, Breach Map, PIPEDA, Legal Icons)"
        },
        {
            "title": "Ethical Considerations",
            "content": "• The Ransom Payment Dilemma: Funding future cybercrime vs. survival\n• Transparency vs. Liability: Delay in transparent communication\n• Downstream Harm: Impact on dealership employees (lost commissions)\n• Consumer Privacy: Protection of sensitive financial PII\n• Systemic Responsibility: Monopoly/Oligopoly duty of care\n• Prioritizing Recovery: Which systems/dealers get restored first?",
            "notes": "The attack presents several profound ethical dilemmas. The most prominent is the ransom payment. While paying the estimated $25 million may have saved dealerships from bankruptcy, it directly funds international cybercrime syndicates—a direct conflict with ethical cybersecurity guidelines like those from the Ransomware Task Force. There was also an ethical failure in transparency; CDK's initial communications were vague, leaving dealers unable to accurately inform their customers. We must also consider the downstream harm to dealership employees whose livelihoods (commissions) were halted, and the duty of care an oligopoly like CDK owes to protect the PII of millions of consumers. [Transition: How do we prevent this going forward?]",
            "image_ref": "Resources #57-61 (Ethics Balance Scale, Ransomware TF Report)"
        },
        {
            "title": "Lessons Learned & Recommendations",
            "content": "• CDK Global: Implement Zero Trust Architecture (ZTA), immutable backups\n• Dealerships: Develop robust manual BCPs, vendor diversification\n• Industry/Auto-ISAC: Standardize API security, require 3rd-party audits\n• Response: Never restore without forensic containment validation\n• Architecture: Eliminate Single Points of Failure in supply chain\n• Testing: Mandate full-scale DR tabletop exercises",
            "notes": "Our recommendations span three tiers. First, SaaS providers like CDK must implement strict Zero Trust Architecture and ensure immutable, air-gapped backups to guarantee recovery. Second, dealerships must realize they cannot entirely outsource their risk. They must develop and regularly test robust, offline manual Business Continuity Plans. Third, the industry—perhaps led by the Auto-ISAC—must demand stricter API security and independent third-party audits of critical vendors. The ultimate lesson learned is structural: the automotive industry must eliminate massive single points of failure and never rush system restoration without complete forensic validation of containment. [Transition: Let's compare this to our previous case study.]",
            "image_ref": "Resources #62-67 (Zero Trust Diagram, 3-2-1 Backup, 3-Tier Recommendation)"
        },
        {
            "title": "Connection to Case Study 1",
            "content": "• Focus: C-Commerce Success & Crisis Readiness\n• Shared Theme 1: High inter-organizational dependency\n• Shared Theme 2: Importance of knowledge exchange during crises\n• Contrast: Internal structural failure vs. external malicious attack\n• Synthesis: Technical DR must pair with human knowledge frameworks\n• Resilience requires both IT availability and collaborative readiness",
            "notes": "Connecting this to our first case study on 'The Influences of Knowledge Exchange on Organizational C-Commerce Success and Crisis Readiness,' we see clear parallels. Both cases highlight the vulnerabilities of highly integrated inter-organizational dependencies. In the first case study, knowledge exchange and collaborative commerce were keys to crisis readiness. In the CDK attack, the lack of secure, out-of-band communication channels severely hampered the crisis response. While Case Study 1 focused more on structural/process crises, CDK is a malicious external attack. The synthesis is clear: technical Disaster Recovery plans are useless without the collaborative human frameworks to execute them during a crisis. [Transition: Let's summarize our key takeaways.]",
            "image_ref": "Resources #68-71 (Venn Diagram, Resilience Model)"
        },
        {
            "title": "Key Takeaways & Discussion",
            "content": "• Takeaway 1: SaaS concentration creates systemic supply chain risk\n• Takeaway 2: Premature recovery exacerbates disaster impacts\n• Takeaway 3: Dealerships require functional offline BCPs\n• Question 1: Should ransom payments be made illegal for critical infrastructure?\n• Question 2: How can dealerships practically diversify their IT dependencies?\n• Question 3: Is the current DMS oligopoly a national security risk?",
            "notes": "To summarize: The CDK attack proves that deep SaaS concentration creates unacceptable systemic supply chain risk. It is a textbook example of how premature recovery attempts can turn an incident into a disaster. Finally, it proves that organizations must maintain functional offline Business Continuity Plans. I'd like to open the floor to discussion with three questions: First, given the impact, should ransom payments be legally banned to disincentivize attacks? Second, how can a mid-sized dealership practically diversify IT when the DMS market is an oligopoly? Third, does this level of software concentration represent a national security risk? [Transition: Thank you for your time.]",
            "image_ref": "Resources #72-75 (Discussion Group, Insight Icons)"
        },
        {
            "title": "References & Q&A",
            "content": "• Tucker, E. (2014). Business Continuity from Preparedness to Recovery\n• IBM Security. (2024). Cost of a Data Breach Report\n• CISA. (2023). #StopRansomware: BlackSuit\n• NIST. (2024). Cybersecurity Framework 2.0\n• Upstream Security. (2024). Global Automotive Cybersecurity Report\n\nQuestions & Peer Critique?",
            "notes": "Here are the primary references used in our analysis, including our course text by Tucker, the IBM Cost of a Data Breach report, and official CISA and NIST documentation. Thank you to our instructor, Sikder Kamruzzaman, and our peers for your attention. We now welcome your questions and constructive critique regarding our BC/DR analysis and presentation delivery. [End of presentation.]",
            "image_ref": "Resources #76-78 (Reference Icon, Q&A Image, Thank You Graphic)"
        }
    ]
    
    for slide_data in slides_data:
        # Use Title and Content layout (index 1) for all slides except Title slide
        layout_idx = 0 if slide_data == slides_data[0] else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_NAVY
        
        # Title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data["title"]
            for p in title_shape.text_frame.paragraphs:
                p.font.name = "Montserrat"
                p.font.color.rgb = WHITE
                p.font.bold = True
        
        # Content
        if "subtitle" in slide_data and layout_idx == 0:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = slide_data["subtitle"]
            for p in subtitle_shape.text_frame.paragraphs:
                p.font.name = "Montserrat"
                p.font.color.rgb = ELECTRIC_BLUE
        elif layout_idx == 1:
            body_shape = slide.placeholders[1]
            body_shape.text = slide_data["content"]
            for p in body_shape.text_frame.paragraphs:
                p.font.name = "Open Sans"
                p.font.color.rgb = WHITE
                p.font.size = Pt(20)
                
            # Add an image placeholder shape at the right side
            left = Inches(8.5)
            top = Inches(2.0)
            width = Inches(4.0)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f"[Image Placeholder]\n{slide_data['image_ref']}"
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 200, 50)
            tf.paragraphs[0].font.size = Pt(14)
            
        # Speaker notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data["notes"]
        
    prs.save("cdk_global_case_study.pptx")

if __name__ == "__main__":
    create_presentation()