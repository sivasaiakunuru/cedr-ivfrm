#!/usr/bin/env python3
"""
Update CEDR PowerPoint with industry-validated Budget & TCO (v2.0)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')

# Find slide 11 (Budget slide - index 10)
slide = prs.slides[10]

# Clear existing content except title
shapes_to_remove = []
for shape in slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text and not text.startswith('Project Budget'):
            shapes_to_remove.append(shape)

# Remove old content shapes
for shape in shapes_to_remove:
    sp = shape.element
    sp.getparent().remove(sp)

# Add note about industry validation
note_text = "Industry-Validated Costs (v2.0) | Sources: AWS, TÜV SÜD, ZCyberSecurity, eSentire, Raspberry Pi"
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.2))
tf = note_box.text_frame
tf.text = note_text
p = tf.paragraphs[0]
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# Add comprehensive budget table - REALISTIC NUMBERS
table_left = Inches(0.5)
table_top = Inches(1.6)
table_width = Inches(12.3)
table_height = Inches(4.2)

# Create table with 8 rows, 4 columns
rows = 8
cols = 4
table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# Set column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.9)
table.columns[3].width = Inches(3.4)

# Header row
headers = ['Category', 'Prototype (1 unit)', 'Pilot (50 veh)', 'Scale (1,000 veh)']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    paragraph.font.size = Pt(11)
    paragraph.alignment = PP_ALIGN.CENTER

# Data rows - REALISTIC INDUSTRY-VALIDATED NUMBERS
data = [
    ['Hardware per Unit', '$293', '$210', '$145'],
    ['NRE/Tooling', '$595', '$12,000', '$53,000'],
    ['Cloud (Annual)', '$135', '$773', '$4,536'],
    ['Compliance', '$0', '$18,000', '$53,000'],
    ['Operations (Annual)', '$0', '$0', '$108,000'],
    ['DevSecOps (Annual)', '$0', '$0', '$6,000'],
    ['Total Year 1', '$1,638', '$41,248', '$369,536'],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        if col_idx == 0:
            paragraph.font.bold = True
        else:
            paragraph.alignment = PP_ALIGN.RIGHT
        
        # Highlight total row
        if row_idx == len(data):
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xe8, 0xf4, 0xe8)
            paragraph.font.bold = True

# Add key metrics below table
metrics_text = """Key Metrics (Industry-Validated):
• Cost per Vehicle: $825 (Pilot) → $370 (Scale Y1) → $189 (Steady-state) • 5-Year TCO: $680K (1,000 vehicles) • AWS Cloud: $4.54/vehicle/year
• ISO 21434: $40,500 | UN R155 CSMS: $25,000 initial + $27K/yr surveillance • Pen Testing: $4,500/test • SOC: $48K/year (8×5 outsourced)"""

metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.4))
tf = metrics_box.text_frame
tf.text = metrics_text
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(9)
    paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save updated presentation
prs.save('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')
print("✅ Updated Slide 11 with Industry-Validated Budget & TCO (v2.0)")
print("✅ Sources: AWS, TÜV SÜD, ZCyberSecurity, eSentire, Raspberry Pi")
print("✅ Presentation saved: CEDR_Final_Presentation.pptx")
