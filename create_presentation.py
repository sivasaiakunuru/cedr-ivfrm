from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xaa, 0xcc, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

# SLIDE 1: Title
add_title_slide(prs, 
    "CEDR: Cybersecurity Event Data Recorder",
    "In-Vehicle Forensic Readiness Module (IV-FRM)\n\nTeam Cyber-Torque | CYB408 Capstone\nSt. Clair College")

# SLIDE 2: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Modern connected vehicles face growing cyber threats (CAN bus attacks, ECU exploits, telematics breaches)",
    "Current Event Data Recorders (EDRs) only capture PHYSICAL crash data — no digital forensic capability",
    "Critical gap: No standardized method to securely log and preserve vehicle cybersecurity incidents",
    "Investigators, automakers, and insurers lack forensically sound digital evidence",
    "Solution needed: Tamper-evident, encrypted, real-time security event logging aligned with ISO/SAE 21434"
])

# SLIDE 3: Solution Overview
add_content_slide(prs, "CEDR Solution Overview", [
    "Embedded system that automatically records and protects vehicle security events",
    "Cryptographic hash chaining (blockchain-style) for tamper evidence",
    "AES-256 encryption for data confidentiality",
    "Real-time event transmission via 4G/5G/WiFi for critical alerts",
    "Cloud backend for fleet correlation and forensic analysis",
    "Web-based investigator dashboard for evidence retrieval",
    "Court-admissible evidence packaging with chain of custody"
])

# SLIDE 4: System Architecture
add_content_slide(prs, "System Architecture", [
    "IN-VEHICLE MODULE (IV-FRM): Event capture → Tamper evidence → Secure storage (SQLite)",
    "EVENT QUEUE & UPLOAD: Immediate upload for critical events, batch for routine",
    "CLOUD BACKEND: Event ingestion → Fleet correlation → Forensic analysis",
    "INVESTIGATOR DASHBOARD: Search, filter, generate reports, verify integrity",
    "COMMUNICATION: 4G/5G/WiFi with HMAC-SHA256 authentication",
    "SECURITY: Fernet (AES-128), PBKDF2 key derivation, RSA-2048 digital signatures"
])

# SLIDE 5: Key Features
add_content_slide(prs, "Key Features & Capabilities", [
    "Tamper-Evident Logging: Blockchain-style hash chain detects any modification",
    "Multi-Layer Cryptography: HMAC (authentication), SHA-256 (integrity), AES-256 (confidentiality)",
    "Real-Time Detection: ~5ms event capture, ~2s critical event upload",
    "13 Attack Scenarios: CAN flooding, message injection, replay, ECU spoofing, MITM, firmware attacks",
    "Fleet Correlation: Cross-vehicle attack pattern analysis",
    "Legal Compliance: Chain of custody, digital signatures, ISO/SAE 21434 alignment"
])

# SLIDE 6: ISO/SAE 21434 Compliance
add_content_slide(prs, "ISO/SAE 21434 Compliance (100%)", [
    "RC-01 Cybersecurity Governance: Audit logs and documentation implemented",
    "RC-02 Risk Management: CVSS scoring integrated for event severity",
    "RC-03 Security by Design: AES-256, HMAC-SHA256, hash chaining",
    "RC-04 Security Validation: Self-tests and integrity verification",
    "RC-05 Security Operations: Real-time event monitoring and alerting",
    "RC-06 Incident Response: Automated alert generation and response",
    "RC-07 Forensic Readiness: Evidence packaging with chain of custody"
])

# SLIDE 7: Attack Simulation Results
add_content_slide(prs, "Attack Simulation Results (13 Scenarios)", [
    "CAN Bus Flooding: ~100ms detection, CRITICAL severity, BLOCKED",
    "Message Injection: ~50ms detection, CRITICAL severity, BLOCKED",
    "Replay Attack: ~150ms detection, HIGH severity, DETECTED",
    "Cellular MITM: ~200ms detection, CRITICAL severity, DETECTED",
    "Firmware Tampering: ~500ms detection, CRITICAL severity, DETECTED",
    "Key Fob Relay: ~100ms detection, HIGH severity, DETECTED",
    "All 13 attacks: 100% detection rate, tamper evidence preserved"
])

# SLIDE 8: Performance Metrics
add_content_slide(prs, "Performance Metrics", [
    "Event Capture Latency: < 10ms target | ~5ms actual ✅",
    "Critical Event Upload: < 5s target | ~2s actual ✅",
    "Chain Verification: < 1s target | ~200ms actual ✅",
    "Evidence Package Creation: < 5s target | ~3s actual ✅",
    "Storage per Event: < 1KB target | ~500B actual ✅",
    "Concurrent Attack Handling: 10+ target | 13 tested ✅"
])

# SLIDE 9: Agile User Stories
add_content_slide(prs, "Agile User Stories", [
    "As a fleet manager, I want real-time security alerts so I can respond to vehicle cyber incidents immediately",
    "As a forensic investigator, I want tamper-evident logs so I can present court-admissible evidence",
    "As an OEM, I want ISO 21434 compliance so I can meet regulatory requirements (UN R155)",
    "SECURITY STORY: As an attacker attempting log tampering, the system shall detect modifications via hash chain breaks",
    "ABUSE STORY: As a malicious insider, attempts to delete logs shall fail due to immutable cloud backup",
    "ACCEPTANCE: All evidence must include chain of custody, digital signatures, and integrity verification"
])

# SLIDE 10: Use Cases
add_content_slide(prs, "Use Cases", [
    "UC-1: Secure OTA Update Logging — Log all firmware updates with cryptographic verification",
    "UC-2: Intrusion Detection & Alert — Real-time IDS alerts while vehicle is in operation",
    "UC-3: Post-Incident Forensics — Investigator retrieves complete event timeline with integrity proof",
    "UC-4: Fleet-Wide Attack Correlation — Identify coordinated attacks across multiple vehicles",
    "UC-5: Insurance Claim Validation — Provide tamper-proof evidence for cybersecurity incidents",
    "UC-6: Compliance Auditing — Generate ISO 21434 work product evidence for auditors"
])

# SLIDE 11: Project Budget (TCO)
add_content_slide(prs, "Project Budget & TCO Estimate", [
    "PROTOTYPE CAPEX: ~$2,500 (Raspberry Pi, CAN shield, sensors, dev tools)",
    "PILOT (50 vehicles): ~$45,000 ($900/vehicle including HSM, cellular, labor)",
    "SCALE (1,000 vehicles): ~$400/vehicle (volume pricing, cloud efficiencies)",
    "OPEX Year 1: ~$15,000 (cloud storage, monitoring, vulnerability management)",
    "OPEX Years 2-3: ~$25,000 (OTA updates, incident response, compliance audits)",
    "Compliance: TARA workshops ($5K), ISO 21434 assessment ($10K), CSMS audit ($15K)"
])

# SLIDE 12: Demo
add_content_slide(prs, "Live Demonstration", [
    "✅ CEDR Cloud Backend running on http://localhost:8080",
    "✅ Event dashboard with real-time security metrics",
    "✅ Tamper-evident logging with blockchain-style hashes",
    "✅ 3 security events currently stored (INTRUSION, UNAUTHORIZED_ACCESS, ANOMALY)",
    "✅ API endpoints for event ingestion and retrieval",
    "DEMO: Investigator dashboard, event verification, forensic report generation"
])

# SLIDE 13: Conclusion
add_content_slide(prs, "Conclusion & Impact", [
    "CEDR fills critical gap in automotive cybersecurity: digital forensics for cyber incidents",
    "Tamper-evident, encrypted, real-time logging meets ISO/SAE 21434 & UN R155 requirements",
    "13 attack scenarios validated with 100% detection rate",
    "Court-admissible evidence packaging with chain of custody",
    "Scalable from prototype to fleet-wide deployment",
    "Future work: HSM integration, ML anomaly detection, V2X security logging"
])

# SLIDE 14: Questions
add_title_slide(prs, "Questions?", "\nTeam Cyber-Torque\nCEDR: Cybersecurity Event Data Recorder\n\nSt. Clair College | CYB408 Capstone")

# Save
prs.save('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')
print("✅ Presentation saved to: /home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx")
