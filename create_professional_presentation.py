#!/usr/bin/env python3
"""
Create a professional cybersecurity business presentation
Style: Corporate, clean, modern - like McKinsey/Big 4 consulting decks
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Color scheme - Professional Cybersecurity
COLORS = {
    'primary_dark': RGBColor(0x0D, 0x1B, 0x2A),      # Deep navy
    'primary': RGBColor(0x1E, 0x3A, 0x5F),           # Navy blue
    'accent': RGBColor(0x00, 0x8B, 0x8B),            # Teal
    'accent_light': RGBColor(0x40, 0xE0, 0xD0),      # Light teal
    'warning': RGBColor(0xFF, 0x6B, 0x35),           # Orange
    'danger': RGBColor(0xE7, 0x4C, 0x3C),            # Red
    'success': RGBColor(0x27, 0xAE, 0x60),           # Green
    'text_dark': RGBColor(0x2C, 0x3E, 0x50),         # Dark gray
    'text_light': RGBColor(0x7F, 0x8C, 0x8D),        # Light gray
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'bg_light': RGBColor(0xF8, 0xF9, 0xFA),          # Off-white
}

def add_title_slide(prs, title, subtitle, presenter="Team Cyber-Torque"):
    """Create a professional title slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background shape - dark navy top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary_dark']
    shape.line.fill.background()
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4), Inches(13.333), Inches(0.15))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(11.7), Inches(1))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['accent_light']
    p.font.name = "Calibri"
    
    # Presenter info
    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(6), Inches(11.7), Inches(1.5))
    tf = info_box.text_frame
    tf.text = presenter
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['text_dark']
    p.font.name = "Calibri"
    
    p = tf.add_paragraph()
    p.text = "CYB408-26W-001 | Automobility Cybersecurity Capstone"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_light']
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "St. Clair College | April 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_light']
    
    return slide

def add_section_divider(prs, section_num, section_title):
    """Create a section divider slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1))
    tf = num_box.text_frame
    tf.text = f"0{section_num}"
    p = tf.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.font.name = "Calibri Light"
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.text = section_title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.5), Inches(3), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    return slide

def add_content_slide(prs, title, bullets, subtitle=None):
    """Create a content slide with bullets"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"
    
    # Subtitle if provided
    start_y = Inches(1.6)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.4))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['accent']
        p.font.name = "Calibri"
        p.font.italic = True
        start_y = Inches(2)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.6), start_y, Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        p.font.name = "Calibri"
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def add_image_slide(prs, title, image_path, caption=None):
    """Create a slide with a full-width image"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"
    
    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.4), width=Inches(12.3))
        print(f"   ✅ Added: {os.path.basename(image_path)}")
    else:
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
        tf = placeholder.text_frame
        tf.text = f"[Image: {os.path.basename(image_path)}]"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['warning']
        print(f"   ⚠️  Missing: {image_path}")
    
    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12), Inches(0.5))
        tf = cap_box.text_frame
        tf.text = caption
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.font.italic = True
        p.font.name = "Calibri"
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content, left_title=None, right_title=None):
    """Create a two-column slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"
    
    # Left column title
    if left_title:
        lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.5))
        tf = lt_box.text_frame
        tf.text = left_title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.font.name = "Calibri"
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_dark']
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.4), Inches(0.02), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['text_light']
    line.line.fill.background()
    
    # Right column title
    if right_title:
        rt_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.5))
        tf = rt_box.text_frame
        tf.text = right_title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.font.name = "Calibri"
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_dark']
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    
    return slide

def add_metrics_slide(prs, title, metrics):
    """Create a metrics/dashboard style slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary_dark']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"
    
    # Metrics boxes
    positions = [
        (Inches(0.5), Inches(1.6)),
        (Inches(4.5), Inches(1.6)),
        (Inches(8.5), Inches(1.6)),
        (Inches(0.5), Inches(4.2)),
        (Inches(4.5), Inches(4.2)),
        (Inches(8.5), Inches(4.2)),
    ]
    
    colors = [COLORS['accent'], COLORS['success'], COLORS['warning'], 
              COLORS['primary'], COLORS['danger'], COLORS['accent_light']]
    
    for i, (metric_name, metric_value, metric_desc) in enumerate(metrics):
        if i >= len(positions):
            break
        x, y = positions[i]
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['bg_light']
        box.line.color.rgb = colors[i % len(colors)]
        box.line.width = Pt(3)
        
        # Metric name
        name_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.5), Inches(0.5))
        tf = name_box.text_frame
        tf.text = metric_name
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_light']
        p.font.name = "Calibri"
        
        # Metric value
        val_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.6), Inches(3.5), Inches(0.8))
        tf = val_box.text_frame
        tf.text = metric_value
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]
        p.font.name = "Calibri Light"
        
        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.5), Inches(3.5), Inches(0.7))
        tf = desc_box.text_frame
        tf.text = metric_desc
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_dark']
        p.font.name = "Calibri"
    
    return slide

def main():
    print("=" * 70)
    print("CREATING PROFESSIONAL CYBERSECURITY PRESENTATION")
    print("Style: Corporate Consulting | Format: Business Professional")
    print("=" * 70)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    print("\n📊 Building presentation slides...\n")
    
    # ===== SLIDE 1: Title =====
    print("1️⃣  Title Slide")
    add_title_slide(prs, 
        "CEDR: Cybersecurity Event Data Recorder",
        "In-Vehicle Forensic Readiness Module for Connected Vehicles",
        "Team Cyber-Torque | CYB408 Capstone")
    
    # ===== SLIDE 2: Executive Summary =====
    print("2️⃣  Executive Summary")
    add_content_slide(prs, "Executive Summary",
        [
            "CEDR detects, records, and preserves cybersecurity events in connected vehicles using tamper-evident blockchain-style hash chaining",
            "Addresses critical gap: OEMs lack forensic readiness for cyber incidents, leading to denied insurance claims and stalled investigations",
            "82% cost advantage over competitors: $485/vehicle vs $800-1,200 for ESCRYPT/Harman solutions",
            "Full ISO/SAE 21434 compliance (7/7 work products) and UN R155/R156 readiness",
            "Production-ready architecture: CM4 Industrial hardware, NXP A71CH HSM, -40°C to +85°C operating range",
            "ML-based anomaly detection reduces false positives by 80% vs rule-based systems"
        ],
        "Investment Ask: $45,000 for Phase 1 Foundation (Months 1-3)")
    
    # ===== SLIDE 3: Problem Statement =====
    print("3️⃣  Problem Statement")
    add_two_column_slide(prs, "The Problem: Forensic Blindness in Connected Vehicles",
        [
            "100+ ECUs per modern vehicle generate terabytes of data",
            "No standardized security event logging exists today",
            "When breaches occur, OEMs cannot prove what happened",
            "Insurance claims denied due to lack of evidence",
            "Regulatory fines for non-compliance (UN R155)",
            "Average breach detection time: 287 days"
        ],
        [
            "Current solutions cost $800-1,200 per vehicle",
            "No tamper-evident logging capabilities",
            "Proprietary systems lack transparency",
            "High false positive rates (20-30%)",
            "Limited to luxury vehicle segments",
            "No open-source alternatives exist"
        ],
        "Market Gap",
        "Current Solutions")
    
    # ===== SECTION DIVIDER =====
    add_section_divider(prs, 1, "Solution Architecture")
    
    # ===== SLIDE 4: System Architecture =====
    print("4️⃣  System Architecture")
    img_path = '/home/siva/.openclaw/workspace/uml/component_diagram_professional.png'
    add_image_slide(prs, "System Architecture", img_path,
                   "Component view showing vehicle edge, multi-cloud backend, and SOC integration")
    
    # ===== SLIDE 5: Deployment Architecture =====
    print("5️⃣  Deployment Architecture")
    img_path = '/home/siva/.openclaw/workspace/uml/deployment_diagram_professional.png'
    add_image_slide(prs, "Deployment Architecture", img_path,
                   "Automotive-grade hardware deployment with secure boot and HSM protection")
    
    # ===== SLIDE 6: Key Features =====
    print("6️⃣  Key Features")
    add_content_slide(prs, "Key Technical Features",
        [
            "Real-time Detection: <500ms event capture with edge ML inference using TensorFlow Lite",
            "Tamper Evidence: Blockchain-style hash chaining ensures court-admissible forensics",
            "Hardware Security: NXP A71CH HSM provides FIPS 140-2 Level 2 key protection",
            "Secure Boot: U-Boot with dm-verity ensures only authentic firmware executes",
            "Multi-Cloud: AWS (primary) + Azure (DR) + GCP (ML) prevents vendor lock-in",
            "Environmental Hardening: IP67 enclosure, -40°C to +85°C operating range",
            "ML Anomaly Detection: <5% false positive rate vs 20-30% industry average",
            "24/7 SOC Integration: eSentire monitoring with automated incident response"
        ])
    
    # ===== SECTION DIVIDER =====
    add_section_divider(prs, 2, "Security Framework")
    
    # ===== SLIDE 7: User Stories =====
    print("7️⃣  User Stories")
    img_path = '/home/siva/.openclaw/workspace/uml/story_pyramid.png'
    if not os.path.exists(img_path):
        img_path = '/home/siva/.openclaw/workspace/visualizations/story_pyramid.png'
    add_image_slide(prs, "Agile User Stories", img_path,
                   "28 user stories: 7 Functional, 7 Security, 4 Non-Functional, 5 Abuse Cases, 5 Countermeasures")
    
    # ===== SLIDE 8: Use Cases =====
    print("8️⃣  Use Cases")
    img_path = '/home/siva/.openclaw/workspace/uml/use_case_diagram.png'
    add_image_slide(prs, "Use Case Diagram", img_path,
                   "6 automotive-specific use cases aligned with ISO/SAE 21434")
    
    # ===== SLIDE 9: Risk Analysis =====
    print("9️⃣  Risk Analysis")
    img_path = '/home/siva/.openclaw/workspace/uml/risk_heatmap.png'
    if not os.path.exists(img_path):
        img_path = '/home/siva/.openclaw/workspace/visualizations/risk_heatmap.png'
    add_image_slide(prs, "Risk Assessment Matrix", img_path,
                   "17 risks analyzed across 5 categories with likelihood/impact scoring")
    
    # ===== SLIDE 10: SWOT =====
    print("🔟 SWOT Analysis")
    img_path = '/home/siva/.openclaw/workspace/uml/swot_diagram.png'
    if not os.path.exists(img_path):
        img_path = '/home/siva/.openclaw/workspace/visualizations/swot_diagram.png'
    add_image_slide(prs, "SWOT Analysis", img_path,
                   "Strategic positioning showing 82% cost advantage and automotive-grade improvements")
    
    # ===== SECTION DIVIDER =====
    add_section_divider(prs, 3, "Business Case")
    
    # ===== SLIDE 11: Financial Metrics =====
    print("1️⃣1️⃣ Financial Metrics")
    add_metrics_slide(prs, "Financial Overview",
        [
            ("Cost per Vehicle", "$485", "Production scale (1,000 units)"),
            ("Cost Advantage", "82%", "vs ESCRYPT/Harman ($800-1,200)"),
            ("5-Year TCO", "$5.98M", "1,000 vehicles, full operations"),
            ("Phase 1 Investment", "$45,000", "Foundation (Months 1-3)"),
            ("ROI", "2.2x", "Risk mitigation value / Investment"),
            ("Break-even", "2,847", "Vehicles for payback")
        ])
    
    # ===== SLIDE 12: Budget Breakdown =====
    print("1️⃣2️⃣ Budget Breakdown")
    add_content_slide(prs, "Project Budget & Total Cost of Ownership",
        [
            "Prototype (5 vehicles): $65,000 total | $1,850 per vehicle | Hardware validation and proof of concept",
            "Pilot (50 vehicles): $385,000 total | $1,850 per vehicle | Field testing with fleet operators",
            "Production (1,000 vehicles): $1,485,000 total | $485 per vehicle | Economies of scale achieved",
            "Year 1 Operations: $1,533,000 | Includes 24/7 SOC ($180K), enterprise SIEM ($48K), personnel (4 FTE)",
            "5-Year TCO: $5,978,000 | Hardware amortized, cloud infrastructure scaling, compliance maintenance",
            "All costs based on commercial rates: Professional engineering salaries, enterprise software licenses, certified penetration testing"
        ],
        "Commercial Pricing - No Academic Discounts Applied")
    
    # ===== SLIDE 13: Improvement Roadmap =====
    print("1️⃣3️⃣ Improvement Roadmap")
    img_path = '/home/siva/.openclaw/workspace/visualizations/improvement_roadmap.png'
    add_image_slide(prs, "4-Phase Improvement Roadmap", img_path,
                   "$1.37M investment over 24 months achieving 90% risk reduction")
    
    # ===== SECTION DIVIDER =====
    add_section_divider(prs, 4, "Compliance & Standards")
    
    # ===== SLIDE 14: Compliance =====
    print("1️⃣4️⃣ Compliance")
    add_two_column_slide(prs, "Regulatory Compliance",
        [
            "RC-01: Cybersecurity Policy ✓",
            "RC-02: Risk Management ✓",
            "RC-03: Security by Design ✓",
            "RC-04: Validation ✓",
            "RC-05: Operations ✓",
            "RC-06: Incident Response ✓",
            "RC-07: Forensic Readiness ✓"
        ],
        [
            "R155: Cyber Security Management System ✓",
            "R156: Software Update Management System ✓",
            "CSMS: Full implementation ✓",
            "SUMS: Secure OTA logging ✓",
            "Incident Response: Automated SOC ✓",
            "Threat Monitoring: ML detection ✓"
        ],
        "ISO/SAE 21434 (7/7 Work Products)",
        "UN R155/R156")
    
    # ===== SLIDE 15: Competitive Analysis =====
    print("1️⃣5️⃣ Competitive Analysis")
    add_two_column_slide(prs, "Competitive Positioning",
        [
            "ESCRYPT: $800-1,200/vehicle",
            "Harman: $600-900/vehicle",
            "Argus: $600-900/vehicle",
            "No open-source alternatives",
            "No tamper-evident solutions"
        ],
        [
            "CEDR: $485/vehicle (82% savings)",
            "Tamper-evident hash chaining",
            "Open-source transparency",
            "ML anomaly detection",
            "Automotive-grade hardware"
        ],
        "Competitors",
        "CEDR Advantage")
    
    # ===== SLIDE 16: Success Metrics =====
    print("1️⃣6️⃣ Success Metrics")
    add_metrics_slide(prs, "Key Performance Indicators",
        [
            ("Detection Latency", "<500ms", "Edge ML inference time"),
            ("False Positive Rate", "<5%", "vs 20-30% industry avg"),
            ("System Uptime", "99.9%", "Multi-cloud SLA"),
            ("Temperature Range", "-40°C to +85°C", "AEC-Q100 Grade 2"),
            ("Attack Coverage", "50+", "Test scenarios"),
            ("Risk Reduction", "63%", "Medium-High to Low-Medium")
        ])
    
    # ===== SLIDE 17: Ask =====
    print("1️⃣7️⃣ The Ask")
    add_content_slide(prs, "Investment Opportunity",
        [
            "We are seeking $45,000 for Phase 1: Foundation (Months 1-3)",
            "Deliverables: 10x CM4 Industrial units, NXP A71CH HSM integration, IP67 enclosures",
            "Milestones: Environmental testing complete, 25 attack scenarios validated, pilot customer LOI",
            "Risk Reduction: 40% improvement (overall risk score 11.3 → 6.8)",
            "Next Steps: Phase 2 hardening ($120K) and Phase 3 production ($285K)",
            "Timeline to Market: 12 months to production readiness"
        ],
        "$45,000 unlocks $1.37M roadmap and positions for Series A")
    
    # ===== SLIDE 18: Team =====
    print("1️⃣8️⃣ Team")
    add_content_slide(prs, "Team Cyber-Torque",
        [
            "St. Clair College | Diploma in Cybersecurity Program",
            "CYB408-26W-001 | Automobility Cybersecurity Capstone",
            "Specialization: Automotive cybersecurity, embedded systems, cloud security",
            "Advisory: Industry mentors from automotive OEMs and cybersecurity firms",
            "Competencies: ISO 21434, UN R155, penetration testing, ML/AI, embedded development",
            "Recognition: St. Clair Cybersecurity Club leadership and CTF competition winners"
        ])
    
    # ===== SLIDE 19: Contact =====
    print("1️⃣9️⃣ Contact")
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary_dark']
    bg.line.fill.background()
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(5.333), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # Thank you
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(1))
    tf = thanks.text_frame
    tf.text = "Thank You"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"
    p.alignment = PP_ALIGN.CENTER
    
    # Questions
    qs = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.3), Inches(0.8))
    tf = qs.text_frame
    tf.text = "Questions & Discussion"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['accent_light']
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5))
    tf = contact.text_frame
    tf.text = "Team Cyber-Torque | CYB408 Capstone"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "St. Clair College | Windsor, ON"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    
    # Save
    output_path = '/home/siva/.openclaw/workspace/CEDR_Professional_Presentation.pptx'
    prs.save(output_path)
    
    print("\n" + "=" * 70)
    print("PRESENTATION CREATION COMPLETE")
    print("=" * 70)
    print(f"\n📊 Total slides: {len(prs.slides)}")
    print(f"💾 Saved to: {output_path}")
    print("\n✅ Professional cybersecurity business presentation ready!")
    print("=" * 70)

if __name__ == '__main__':
    main()
