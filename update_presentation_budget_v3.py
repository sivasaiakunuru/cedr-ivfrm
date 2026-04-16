#!/usr/bin/env python3
"""
Update CEDR PowerPoint with Business Implementation Budget (v3.0)
Commercial pricing - no student discounts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing presentation
prs = Presentation('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')

# Find slide 11 (Budget slide - index 10)
slide = prs.slides[10]

# Clear existing content except title
shapes_to_remove = []
for shape in slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text and not text.startswith('Project Budget'):
            shapes_to_remove.append(shape)

# Remove old content shapes
for shape in shapes_to_remove:
    sp = shape.element
    sp.getparent().remove(sp)

# Add note about commercial implementation
note_text = "Business Implementation (Commercial Pricing) | Market rates for professional deployment"
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.2))
tf = note_box.text_frame
tf.text = note_text
p = tf.paragraphs[0]
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# Add comprehensive budget table - COMMERCIAL PRICING
table_left = Inches(0.5)
table_top = Inches(1.6)
table_width = Inches(12.3)
table_height = Inches(4.2)

# Create table with 8 rows, 4 columns
rows = 8
cols = 4
table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# Set column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.9)
table.columns[3].width = Inches(3.4)

# Header row
headers = ['Category', 'Prototype (POC)', 'Pilot (50 veh)', 'Production (1,000)']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    paragraph.font.size = Pt(11)
    paragraph.alignment = PP_ALIGN.CENTER

# Data rows - COMMERCIAL PRICING (NO STUDENT DISCOUNTS)
data = [
    ['Engineering Labor', '$43,000', '$110,000', '$615,000'],
    ['Hardware per Unit', '$771', '$290', '$145'],
    ['NRE/Tooling', '$20,885', '$50,500', '$140,000'],
    ['Cloud (Annual)', '$625', '$2,286', '$12,096'],
    ['Enterprise Software', '$4,860', '$57,600', '$174,000'],
    ['Compliance & Security', '$0', '$78,000', '$202,000'],
    ['Total Year 1', '$65,000', '$385,000', '$1,485,000'],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        if col_idx == 0:
            paragraph.font.bold = True
        else:
            paragraph.alignment = PP_ALIGN.RIGHT
        
        # Highlight total row
        if row_idx == len(data):
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xe8, 0xf4, 0xe8)
            paragraph.font.bold = True

# Add key metrics below table
metrics_text = """Key Metrics (Commercial Implementation):
• Cost per Vehicle: $1,850 (Pilot) → $485 (Production) → $405 (Steady-state) • 5-Year TCO: $5.43M (1,000 vehicles)
• Cloud: $12.10/vehicle/year • SOC: $180K/year (24/7 outsourced) • ISO 21434: $53K | UN R155 CSMS: $67K
• Payback Period: 9.9 years (via risk mitigation) • Competitive vs. commercial alternatives ($600-1,200/vehicle)"""

metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.4))
tf = metrics_box.text_frame
tf.text = metrics_text
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(9)
    paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save updated presentation
prs.save('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')
print("✅ Updated Slide 11 with Business Implementation Budget (v3.0)")
print("✅ Commercial pricing - market rates for professional deployment")
print("✅ Presentation saved: CEDR_Final_Presentation.pptx")
