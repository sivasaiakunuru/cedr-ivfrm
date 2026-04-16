#!/usr/bin/env python3
"""
Convert HTML presentation to PowerPoint (.pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme (matching HTML dark theme)
DARK_BG = RGBColor(0x0A, 0x25, 0x40)  # #0a2540
ACCENT = RGBColor(0x00, 0xD4, 0xAA)   # #00d4aa
PURPLE = RGBColor(0x63, 0x5B, 0xFF)   # #635bff
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

def add_title_slide(prs, title, subtitle, badges=None):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Logo box
    logo = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.5), Inches(2.3), Inches(2.3)
    )
    logo.fill.gradient()
    logo.fill.gradient_stops[0].color.rgb = PURPLE
    logo.fill.gradient_stops[1].color.rgb = ACCENT
    logo.line.fill.background()
    
    logo_text = logo.text_frame
    logo_text.text = "CEDR"
    p = logo_text.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.3), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.3), Inches(0.6))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Badges
    if badges:
        badge_y = Inches(5.8)
        badge_spacing = Inches(3.5)
        start_x = (prs.slide_width - len(badges) * badge_spacing) / 2 + Inches(1)
        
        for i, badge in enumerate(badges):
            badge_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                start_x + i * badge_spacing - Inches(1.5), badge_y,
                Inches(2.5), Inches(0.4)
            )
            badge_shape.fill.gradient()
            badge_shape.fill.gradient_stops[0].color.rgb = PURPLE
            badge_shape.fill.gradient_stops[1].color.rgb = ACCENT
            badge_shape.line.fill.background()
            
            tf = badge_shape.text_frame
            tf.text = badge
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4))
    tf = footer.text_frame
    tf.text = "Team Cyber-Torque | CYB408 Capstone | St. Clair College"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """Add section header slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, two_column=False):
    """Add content slide with bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Content
    if two_column and len(bullets) >= 2:
        # Split bullets into two columns
        mid = len(bullets) // 2
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(5.8))
        tf = left_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets[:mid]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.space_before = Pt(12)
            p.level = 0
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(6), Inches(5.8))
        tf = right_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets[mid:]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.space_before = Pt(12)
            p.level = 0
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.space_before = Pt(16)
            p.level = 0
    
    return slide

def add_stats_slide(prs, title, stats):
    """Add slide with statistics cards"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Stats cards
    num_stats = len(stats)
    card_width = Inches(3.5)
    spacing = Inches(4)
    start_x = (prs.slide_width - (num_stats * spacing - Inches(0.5))) / 2
    
    for i, (number, label) in enumerate(stats):
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_x + i * spacing, Inches(2),
            card_width, Inches(3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
        card.line.color.rgb = RGBColor(0x33, 0x55, 0x77)
        
        # Number
        num_box = slide.shapes.add_textbox(
            start_x + i * spacing + Inches(0.2), Inches(2.5),
            Inches(3.1), Inches(1.2)
        )
        tf = num_box.text_frame
        tf.text = number
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(
            start_x + i * spacing + Inches(0.2), Inches(3.6),
            Inches(3.1), Inches(1)
        )
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add slide with table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.3),
        Inches(12.3), Inches(0.6 * num_rows)
    ).table
    
    # Set column widths (convert to integer EMUs)
    from pptx.util import Emu
    col_width = int(Inches(12.3).emu / num_cols)
    for col in table.columns:
        col.width = Emu(col_width)
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
    
    return slide

def add_timeline_slide(prs, title, phases):
    """Add roadmap/timeline slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Timeline line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(2),
        Inches(0.05), Inches(5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    # Phase items
    for i, (phase_name, details, cost, risk) in enumerate(phases):
        y_pos = Inches(1.5 + i * 1.4)
        
        # Dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.85), y_pos, Inches(0.3), Inches(0.3)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        
        # Phase box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), y_pos - Inches(0.1),
            Inches(11), Inches(1.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
        box.line.color.rgb = RGBColor(0x33, 0x55, 0x77)
        
        # Phase name
        name_box = slide.shapes.add_textbox(Inches(1.6), y_pos + Inches(0.05), Inches(7), Inches(0.4))
        tf = name_box.text_frame
        tf.text = phase_name
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        
        # Details
        detail_box = slide.shapes.add_textbox(Inches(1.6), y_pos + Inches(0.4), Inches(6), Inches(0.5))
        tf = detail_box.text_frame
        tf.text = details
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        
        # Cost and risk
        cost_box = slide.shapes.add_textbox(Inches(9), y_pos + Inches(0.15), Inches(3), Inches(0.6))
        tf = cost_box.text_frame
        tf.text = f"{cost} | Risk: {risk}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT
    
    return slide

def add_closing_slide(prs, title, subtitle, contact_info):
    """Add closing/Q&A slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.6))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Badges
    badges = ["MIT Licensed", "Open Source", "Production Ready"]
    badge_y = Inches(4.6)
    badge_spacing = Inches(3.5)
    start_x = (prs.slide_width - len(badges) * badge_spacing) / 2 + Inches(1.5)
    
    for i, badge in enumerate(badges):
        badge_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_x + i * badge_spacing - Inches(1.5), badge_y,
            Inches(2.5), Inches(0.4)
        )
        badge_shape.fill.gradient()
        # Note: Gradient setup simplified for compatibility
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = PURPLE
        badge_shape.line.fill.background()
        
        tf = badge_shape.text_frame
        tf.text = badge
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Contact box
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.4),
        Inches(6.3), Inches(1.5)
    )
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    contact_box.line.color.rgb = RGBColor(0x33, 0x55, 0x77)
    
    tf = contact_box.text_frame
    tf.text = contact_info
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== BUILD PRESENTATION ====================

print("Building CEDR Presentation...")

# Slide 1: Title
add_title_slide(prs, 
    "Cybersecurity Event Data Recorder",
    "In-Vehicle Forensic Readiness Module",
    ["ISO/SAE 21434 Compliant", "UN R155/R156 Ready", "ML-Powered"]
)

# Slide 2: The Problem
add_content_slide(prs, "The Critical Challenge", [
    "100+ ECUs: Modern vehicles contain 100+ ECUs generating terabytes of data with no standardized security event logging.",
    "287 Days: Average breach detection time in automotive systems, leaving attackers with extended dwell time.",
    "40% Denial: Insurance claim denial rate due to lack of tamper-evident forensic evidence after cyber incidents.",
    "€30M Fines: Regulatory fines under UN R155 for non-compliance with cybersecurity management systems."
])

# Slide 3: Solution Overview
add_stats_slide(prs, "CEDR Solution Overview", [
    ("$485", "Per Vehicle Cost"),
    ("82%", "Cost Advantage vs Competitors"),
    ("<500ms", "Detection Latency")
])

# Slide 4: System Architecture
add_content_slide(prs, "System Architecture - Vehicle Edge Layer", [
    "CAN Bus Interface: Real-time vehicle network monitoring with <10ms latency",
    "GPS/GNSS Receiver: Location and timestamp correlation for forensic evidence",
    "4G/LTE Cellular: Secure cloud connectivity via encrypted TLS 1.3",
    "CEDR Core Module: Raspberry Pi CM4 Industrial with TensorFlow Lite ML inference",
    "NXP A71CH HSM: Secure key storage and crypto acceleration (FIPS 140-2 Level 2)",
    "Tamper Detection: Enclosure breach sensors and temperature monitoring"
])

# Slide 5: Cloud Infrastructure
add_content_slide(prs, "Cloud Infrastructure", [
    "Multi-Cloud Architecture: AWS (Primary) + Azure (Backup/DR) + GCP (ML Training)",
    "Kubernetes Cluster: Containerized microservices for scalability and resilience",
    "API Gateway: Kong/AWS for secure API management and rate limiting",
    "Event Processor: Real-time stream processing with Apache Kafka",
    "ML Inference: Vertex AI for anomaly detection model serving",
    "PostgreSQL + InfluxDB: Relational and time-series data storage",
    "Blockchain Audit Log: Immutable evidence chain for legal proceedings"
])

# Slide 6: Security Architecture
add_content_slide(prs, "Security Architecture", [
    "AES-256-GCM: Encryption for data at rest with hardware acceleration",
    "SHA-256 Hash Chaining: Tamper-evident logging with cryptographic integrity",
    "TLS 1.3: Latest transport security for data in transit",
    "NXP A71CH Secure Element: Hardware-based key protection",
    "Physical Security: Enclosure breach detection and anti-tamper mesh",
    "Secure Boot: Verified boot chain with signed firmware",
    "Zero-Trust Architecture: Continuous verification and micro-segmentation"
], two_column=True)

# Slide 7: Agile User Stories
add_content_slide(prs, "Agile User Stories", [
    "28 comprehensive user stories covering functional, security, and abuse cases",
    "Functional Stories: Real-time alerts, forensic search, fleet correlation, OTA logging",
    "Security Stories: Tamper evidence, encryption, MFA, RBAC, HSM integration",
    "Abuse Cases: Attacker modifying logs, MITM attacks, insider threats",
    "Countermeasures: Immutable audit trail, certificate pinning, DLP monitoring",
    "Traceability Matrix: Each story maps to requirements → threats → tests → components"
])

# Slide 8: Use Cases
add_content_slide(prs, "Automotive Use Cases", [
    "Secure Vehicle Access: Biometric authentication with forensic logging of all attempts",
    "OTA Update Authentication: Cryptographic verification with rollback protection",
    "IDS Alert Management: Real-time intrusion detection with automated response",
    "Secure ECU-to-Cloud: End-to-end encrypted telemetry with mutual authentication",
    "Driver Identity Verification: Continuous authentication with behavioral biometrics",
    "Misuse Case: Firmware flashing attacks detected and logged for investigation"
], two_column=True)

# Slide 9: UML Diagrams
add_content_slide(prs, "UML Diagram Suite", [
    "Use Case Diagrams: Driver, Mechanic, Cloud Service, Attacker actors",
    "Sequence Diagrams: Message flows Vehicle→Cloud, ECU→IDS→Alert Service",
    "Component Diagrams: CAN, Automotive Ethernet, ECU layout, HSMs, Cloud",
    "Deployment Diagrams: Physical architecture across edge and cloud regions",
    "Activity Diagrams: Incident detection → Logging → Alert escalation",
    "Class Diagrams: SecurityEvent, HashChain, Evidence, Alert data models"
])

# Slide 10: Risk Analysis
add_content_slide(prs, "Risk Analysis & SWOT", [
    "Supply Chain Risk: HIGH - Mitigated via SBOM and vendor validation",
    "Crypto Agility: HIGH - Post-quantum cryptography migration plan",
    "Scalability: MEDIUM - Auto-scaling cloud infrastructure",
    "Integration: MEDIUM - Standard APIs and HAL abstraction layer",
    "Strengths: 82% cost advantage, ISO 21434 compliant, <500ms latency",
    "Opportunities: UN R155 mandate, insurance partnerships, OEM integration"
], two_column=True)

# Slide 11: Project Budget
add_table_slide(prs, "Project Budget & TCO", 
    ["Phase", "Investment", "Timeline", "Risk Reduction", "Deliverables"],
    [
        ["Phase 1: Foundation", "$45,000", "Months 1-3", "40%", "Prototype, basic logging"],
        ["Phase 2: Hardening", "$120,000", "Months 4-6", "65%", "ML detection, encryption"],
        ["Phase 3: Production", "$285,000", "Months 7-12", "80%", "AEC-Q100, compliance audit"],
        ["Phase 4: Scale", "$850,000", "Months 13-24", "90%", "1,000 vehicle deployment"],
        ["TOTAL", "$1.37M", "24 months", "-", "Full production system"]
    ]
)

# Slide 12: Hardware BOM
add_table_slide(prs, "Hardware & BOM",
    ["Component", "Prototype", "Production (1K+)"],
    [
        ["Raspberry Pi CM4 / Custom ARM", "$85", "$65"],
        ["NXP A71CH Secure Element", "$12", "$8"],
        ["CAN Bus / Ethernet PHY", "$25", "$15"],
        ["4G/LTE Cellular Module", "$45", "$32"],
        ["GPS/GNSS Receiver", "$18", "Integrated"],
        ["Enclosure & Sensors", "$35", "$25"],
        ["Assembly & Testing", "-", "$25"],
        ["Certification (amortized)", "-", "$40"],
        ["TOTAL UNIT COST", "$220", "$185"]
    ]
)

# Slide 13: Compliance
add_content_slide(prs, "Regulatory Compliance", [
    "ISO/SAE 21434: Complete coverage (7/7 work products)",
    "UN R155: Cybersecurity Management System (CSMS) compliant",
    "UN R156: Software Update Management System (SUMS) ready",
    "AEC-Q100: Automotive qualification (Phase 3)",
    "CSMS: Organizational governance, risk management, incident response",
    "SUMS: Software update process, verification, user notification",
    "Forensic evidence satisfies UN R155 incident response obligations"
])

# Slide 14: DevSecOps
add_content_slide(prs, "DevSecOps & Tooling", [
    "CI/CD Pipeline: GitHub Actions, Docker, Kubernetes, Terraform",
    "Security Testing: SAST (SonarQube), DAST (OWASP ZAP), SCA (Snyk), Fuzzing (AFL)",
    "Artifact Management: SBOM generation, firmware signing, secret management",
    "Security Champion Model: Embedded security stakeholders in dev teams",
    "Shift-Left Integration: Threat modeling in sprint planning",
    "Abuse case identification before architecture hardens"
])

# Slide 15: Cloud Costs
add_table_slide(prs, "Cloud Infrastructure Costs (AWS)",
    ["Service", "Monthly Cost"],
    [
        ["EC2 (Kubernetes nodes)", "$850"],
        ["S3 Storage (30-day retention)", "$420"],
        ["Data Transfer (Egress)", "$380"],
        ["CloudWatch (Monitoring)", "$180"],
        ["KMS (Key Management)", "$120"],
        ["Secrets Manager", "$80"],
        ["Support Plan (Business)", "$250"],
        ["TOTAL MONTHLY", "$2,280"],
        ["ANNUAL COST", "$27,360"]
    ]
)

# Slide 16: Roadmap
add_timeline_slide(prs, "Development Roadmap", [
    ("Phase 1: Foundation", "Hardware prototyping, basic CAN logging, cloud connectivity", "$45,000", "40%"),
    ("Phase 2: Hardening", "ML model integration, cryptographic protections, tamper detection", "$120,000", "65%"),
    ("Phase 3: Production", "AEC-Q100 qualification, compliance audits, pilot fleet (50 vehicles)", "$285,000", "80%"),
    ("Phase 4: Scale", "Mass production, 1,000 vehicle deployment, continuous monitoring", "$850,000", "90%")
])

# Slide 17: Team
add_content_slide(prs, "Team Cyber-Torque", [
    "Course: CYB408-26W-001 Automobility Cybersecurity CAP",
    "Institution: St. Clair College, Chatham, Ontario",
    "Project Type: Capstone Project 2024-2025",
    "Advisors: Industry consultants and academic reviewers",
    "Repository: github.com/sivasaiakunuru/cedr-ivfrm",
    "Status: Production ready for pilot deployment"
])

# Slide 18: Closing
add_closing_slide(prs,
    "Questions?",
    "CEDR-IVFRM Project - Cybersecurity Event Data Recorder for Connected Vehicles",
    "Contact: Team Cyber-Torque\nEmail: capstone@cybertorque.ca\nGitHub: github.com/sivasaiakunuru/cedr-ivfrm"
)

# Save presentation
output_file = "CEDR_Presentation_Web_Converted.pptx"
prs.save(output_file)
print(f"✅ Presentation saved: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
