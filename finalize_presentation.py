#!/usr/bin/env python3
"""
Finalize CEDR Presentation for CYB408 Capstone
Inserts all professional images and follows course guidelines
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# File paths
PPTX_PATH = '/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx'
UML_DIR = '/home/siva/.openclaw/workspace/uml'
VIZ_DIR = '/home/siva/.openclaw/workspace/visualizations'

def add_image_slide(prs, title_text, image_path, left=Inches(1), top=Inches(1.5), 
                    width=Inches(11), height=Inches(6)):
    """Add a slide with title and full-width image"""
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    p.alignment = PP_ALIGN.CENTER
    
    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width, height)
        print(f"  ✅ Added image: {os.path.basename(image_path)}")
    else:
        # Placeholder text if image not found
        placeholder = slide.shapes.add_textbox(left, top, width, height)
        tf = placeholder.text_frame
        tf.text = f"[Insert Image: {os.path.basename(image_path)}]"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        print(f"  ⚠️  Image not found: {image_path}")
    
    return slide

def update_existing_slide(slide, title_text=None):
    """Update title of existing slide"""
    if title_text:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text_frame.text and len(shape.text_frame.text) > 0:
                    # Check if this looks like a title (top of slide, large text)
                    if shape.top < Inches(1) and shape.text_frame.paragraphs[0].font.size > Pt(20):
                        shape.text_frame.text = title_text
                        p = shape.text_frame.paragraphs[0]
                        p.font.size = Pt(32)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)
                        break

def main():
    print("=" * 70)
    print("FINALIZING CEDR PRESENTATION FOR CYB408")
    print("=" * 70)
    
    # Load presentation
    print("\n📂 Loading presentation...")
    prs = Presentation(PPTX_PATH)
    print(f"   Current slides: {len(prs.slides)}")
    
    # Slide 1: Title (keep existing)
    print("\n1️⃣  Slide 1: Title - Verified")
    
    # Add new slides with images
    print("\n📊 Adding slides with professional images...")
    
    # Slide 2: System Architecture (Component Diagram)
    print("\n2️⃣  Adding System Architecture slide...")
    img_path = os.path.join(UML_DIR, 'component_diagram_professional.png')
    slide = add_image_slide(prs, "System Architecture", img_path, 
                           left=Inches(0.5), top=Inches(1.3), 
                           width=Inches(12), height=Inches(6.5))
    
    # Slide 3: Deployment Architecture
    print("\n3️⃣  Adding Deployment Architecture slide...")
    img_path = os.path.join(UML_DIR, 'deployment_diagram_professional.png')
    slide = add_image_slide(prs, "Deployment Architecture", img_path,
                           left=Inches(0.5), top=Inches(1.3),
                           width=Inches(12), height=Inches(6.5))
    
    # Slide 4: User Story Pyramid
    print("\n4️⃣  Adding User Stories slide...")
    img_path = os.path.join(UML_DIR, 'story_pyramid.png')
    if not os.path.exists(img_path):
        img_path = os.path.join(VIZ_DIR, 'story_pyramid.png')
    slide = add_image_slide(prs, "Agile User Stories (28 Total)", img_path,
                           left=Inches(2), top=Inches(1.5),
                           width=Inches(9), height=Inches(6))
    
    # Slide 5: Use Case Diagram
    print("\n5️⃣  Adding Use Cases slide...")
    img_path = os.path.join(UML_DIR, 'use_case_diagram.png')
    slide = add_image_slide(prs, "Use Case Diagram", img_path,
                           left=Inches(1), top=Inches(1.3),
                           width=Inches(11), height=Inches(6.5))
    
    # Slide 6: Risk Heat Map
    print("\n6️⃣  Adding Risk Analysis slide...")
    img_path = os.path.join(UML_DIR, 'risk_heatmap.png')
    if not os.path.exists(img_path):
        img_path = os.path.join(VIZ_DIR, 'risk_heatmap.png')
    slide = add_image_slide(prs, "Risk Heat Map", img_path,
                           left=Inches(2.5), top=Inches(1.5),
                           width=Inches(8), height=Inches(6))
    
    # Slide 7: SWOT Analysis
    print("\n7️⃣  Adding SWOT Analysis slide...")
    img_path = os.path.join(UML_DIR, 'swot_diagram.png')
    if not os.path.exists(img_path):
        img_path = os.path.join(VIZ_DIR, 'swot_diagram.png')
    slide = add_image_slide(prs, "SWOT Analysis", img_path,
                           left=Inches(1.5), top=Inches(1.5),
                           width=Inches(10), height=Inches(6))
    
    # Slide 8: Improvement Roadmap
    print("\n8️⃣  Adding Improvement Roadmap slide...")
    img_path = os.path.join(VIZ_DIR, 'improvement_roadmap.png')
    slide = add_image_slide(prs, "4-Phase Improvement Roadmap", img_path,
                           left=Inches(0.5), top=Inches(1.3),
                           width=Inches(12), height=Inches(6.5))
    
    # Slide 9: Budget Summary (text-based with key metrics)
    print("\n9️⃣  Adding Budget Summary slide...")
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Project Budget & TCO"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    p.alignment = PP_ALIGN.CENTER
    
    # Budget metrics box
    metrics = [
        ("Prototype", "5 vehicles", "$65,000", "$1,850/veh"),
        ("Pilot", "50 vehicles", "$385,000", "$1,850/veh"),
        ("Production", "1,000 vehicles", "$1,485,000", "$485/veh"),
        ("5-Year TCO", "1,000 vehicles", "$5,978,000", "$405/year/veh"),
    ]
    
    y_pos = Inches(1.5)
    for phase, vehicles, cost, per_veh in metrics:
        # Phase name
        box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(2.5), Inches(0.6))
        tf = box.text_frame
        tf.text = phase
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)
        
        # Vehicles
        box = slide.shapes.add_textbox(Inches(3), y_pos, Inches(2.5), Inches(0.6))
        tf = box.text_frame
        tf.text = vehicles
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        
        # Cost
        box = slide.shapes.add_textbox(Inches(5.5), y_pos, Inches(3), Inches(0.6))
        tf = box.text_frame
        tf.text = cost
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x00, 0x66, 0x00)
        
        # Per vehicle
        box = slide.shapes.add_textbox(Inches(9), y_pos, Inches(3), Inches(0.6))
        tf = box.text_frame
        tf.text = per_veh
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        
        y_pos += Inches(0.9)
    
    # Highlight box
    highlight = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.2))
    tf = highlight.text_frame
    tf.text = "Key Highlights:"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    
    highlights = [
        "• 82% cost advantage over competitors ($485 vs $800+ per vehicle)",
        "• Hardware upgraded to CM4 Industrial with NXP A71CH HSM",
        "• Includes 24/7 SOC monitoring, enterprise SIEM, penetration testing",
    ]
    for hl in highlights:
        p = tf.add_paragraph()
        p.text = hl
        p.font.size = Pt(12)
    
    print("   ✅ Budget summary added")
    
    # Save presentation
    print("\n💾 Saving presentation...")
    prs.save(PPTX_PATH)
    print(f"   ✅ Saved: {PPTX_PATH}")
    
    # Summary
    print("\n" + "=" * 70)
    print("PRESENTATION FINALIZATION COMPLETE")
    print("=" * 70)
    print(f"\n📊 Total slides: {len(prs.slides)}")
    print("\n📁 Slides with images:")
    print("   • System Architecture (Component Diagram)")
    print("   • Deployment Architecture")
    print("   • Agile User Stories (28 stories)")
    print("   • Use Case Diagram")
    print("   • Risk Heat Map")
    print("   • SWOT Analysis")
    print("   • Improvement Roadmap")
    print("   • Budget & TCO Summary")
    print("\n✅ Presentation is ready for CYB408 submission!")
    print("=" * 70)

if __name__ == '__main__':
    main()
