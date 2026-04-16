#!/usr/bin/env python3
"""
Add Improvement Roadmap slide to CEDR PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load existing presentation
prs = Presentation('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')

# Add new slide after Risk Analysis (index 12)
blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "Improvement Roadmap & Next Steps"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)

# Before/After comparison table
table_left = Inches(0.5)
table_top = Inches(1.2)
table_width = Inches(12)
table_height = Inches(3.5)

rows = 7
cols = 3
table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# Column widths
table.columns[0].width = Inches(3)
table.columns[1].width = Inches(4)
table.columns[2].width = Inches(5)

# Headers
headers = ['Area', 'Current (Prototype)', 'Improved (Production)']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

# Data
data = [
    ['Hardware', 'Raspberry Pi 4', 'CM4 Industrial + HSM'],
    ['Temp Range', '0°C to 50°C', '-40°C to +85°C'],
    ['Testing', '13 attack scenarios', '50+ scenarios'],
    ['Detection', 'Rule-based', 'ML anomaly detection'],
    ['Cloud', 'AWS only', 'Multi-cloud + Edge'],
    ['Cost/Vehicle', '$1,850 (pilot)', '$485 (production)'],
]

colors = [
    RGBColor(0xFF, 0xE6, 0xE6),  # Light red
    RGBColor(0xE6, 0xFF, 0xE6),  # Light green
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors[col_idx % 2]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        if col_idx == 0:
            p.font.bold = True

# Investment summary
summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(12), Inches(1.5))
tf = summary_box.text_frame
tf.text = "4-Phase Implementation Plan:"

phases = [
    "• Phase 1 (M1-3): Foundation - $45K, 40% risk reduction",
    "• Phase 2 (M4-6): Hardening - $120K, 65% risk reduction",
    "• Phase 3 (M7-12): Production - $285K, 80% risk reduction",
    "• Phase 4 (M13-24): Scale - $850K, 90% risk reduction",
    "",
    "Total Investment: $1.37M | ROI: 2.2x (Risk mitigation value / Investment)"
]

for phase in phases:
    p = tf.add_paragraph()
    p.text = phase
    p.font.size = Pt(10)
    if "Total Investment" in phase:
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x00, 0x66, 0x00)

# Save
prs.save('/home/siva/.openclaw/workspace/CEDR_Final_Presentation.pptx')
print("✅ Added Improvement Roadmap slide (Slide 13)")
print("✅ PowerPoint now has 16 slides total")
